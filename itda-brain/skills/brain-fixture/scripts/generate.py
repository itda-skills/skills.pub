#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""brain-fixture generate — 원장(ledger)에서 데이터셋 폴더를 결정론 렌더.

사용:
  python3 scripts/generate.py <ledger.json> --out <폴더>   # macOS/Linux
  py -3 scripts/generate.py <ledger.json> --out <폴더>      # Windows

- 원장 documents[] 전건을 유형별로 렌더(docx·xlsx·pptx·pdf·txt·csv + broken·lock).
- os.utime 으로 파일 mtime = 원장 내부 날짜(internal_date).
- 출력 폴더가 비어있지 않으면 exit 2(기존 데이터 덮어쓰기 금지, no-silent-fallback).
- 원장 스키마 위반은 exit 2 명시 에러(어느 필드가 왜).
값 수준 결정론: 같은 원장 → 같은 값·구조(바이트 동일성은 비보장).
"""
from __future__ import annotations

import argparse
import os
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import bf_common as bf  # noqa: E402


def _prepare_out(out: Path) -> None:
    if out.exists():
        if not out.is_dir():
            raise bf.BFError(f"출력 경로가 폴더가 아닙니다: {out}")
        if any(out.iterdir()):
            raise bf.BFError(
                f"출력 폴더가 비어있지 않습니다: {out}\n"
                "기존 데이터를 덮어쓰지 않습니다 — 빈 폴더 또는 새 경로를 지정하세요."
            )
    else:
        out.mkdir(parents=True, exist_ok=True)


def _abspath(out: Path, rel: str) -> Path:
    # 방어 심층(FIX-1) — 스키마 검증이 이미 경로 탈출을 거부하지만, 파일 쓰기 직전 정규화 후
    # 출력 루트 밖이면 차단해 폴더 밖 파일 덮어쓰기를 원천 봉쇄한다.
    p = bf.safe_join(out, rel)
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


def _set_mtime(path: Path, internal_date: str) -> None:
    ts = bf.mtime_ts(internal_date)
    os.utime(path, (ts, ts))


# ---------------------------------------------------------------- 렌더러

def _render_docx(path: Path, doc: dict) -> None:
    from docx import Document

    d = Document()
    d.add_heading(doc["title"], level=0)
    for blk in doc["blocks"]:
        kind = blk["kind"]
        if kind == "p":
            d.add_paragraph(blk["text"])
        elif kind == "h":
            d.add_heading(blk["text"], level=1)
        elif kind == "table":
            rows = blk["rows"]
            if not rows:
                continue
            t = d.add_table(rows=len(rows), cols=len(rows[0]))
            t.style = "Table Grid"
            for i, row in enumerate(rows):
                for j, cell in enumerate(row):
                    t.cell(i, j).text = bf._cell_text(cell)
    d.save(str(path))


def _render_xlsx(path: Path, doc: dict) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    for name, rows in doc["sheets"].items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(list(row))
    wb.save(str(path))


def _render_pptx(path: Path, doc: dict) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    for idx, sl in enumerate(doc["slides"]):
        bullets = [str(b) for b in sl["bullets"]]
        if idx == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = sl["title"]
            if bullets:
                slide.placeholders[1].text = "\n".join(bullets)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sl["title"]
            body = slide.placeholders[1].text_frame
            body.clear()
            for i, b in enumerate(bullets):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = b
                para.font.size = Pt(18)
    prs.save(str(path))


_PDF_FONT_READY = False


def _ensure_pdf_font() -> str:
    global _PDF_FONT_READY
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    font = "HYSMyeongJo-Medium"  # reportlab 내장 CJK 폰트(한글 텍스트 레이어)
    if not _PDF_FONT_READY:
        pdfmetrics.registerFont(UnicodeCIDFont(font))
        _PDF_FONT_READY = True
    return font


def _wrap(line: str, width: int = 40) -> list[str]:
    if len(line) <= width:
        return [line]
    return [line[i : i + width] for i in range(0, len(line), width)]


def _render_pdf(path: Path, doc: dict) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    font = _ensure_pdf_font()
    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4
    y = height - 60

    def _draw(text: str, size: int) -> None:
        nonlocal y
        for seg in _wrap(text):
            if y < 60:
                c.showPage()
                c.setFont(font, size)
                y = height - 60
            c.setFont(font, size)
            c.drawString(50, y, seg)
            y -= size + 8

    _draw(doc["title"], 15)
    y -= 6
    for ln in doc["lines"]:
        _draw(str(ln), 11)
    c.save()


def _render_txt(path: Path, doc: dict) -> None:
    path.write_text(doc["content"].strip() + "\n", encoding="utf-8")


def _render_csv(path: Path, doc: dict) -> None:
    import csv as csvmod

    with path.open("w", encoding="utf-8", newline="") as f:
        w = csvmod.writer(f)
        for row in doc["rows"]:
            w.writerow([bf._cell_text(c) for c in row])


def _render_broken(path: Path, doc: dict) -> None:
    """손상된 zip(오피스 문서 흉내) — 시그니처만 있고 내부가 깨진 바이트열.

    페이로드는 note+경로에서 파생한 고정 시드 기반 의사난수(결정론 — Date/무시드 random 금지,
    ASCII 마커 없음). is_zipfile=False 를 유지해 verify ④축이 손상으로 판정한다.
    """
    import hashlib
    import random

    seed_src = (str(doc.get("note", "")) + "|" + doc["path"]).encode("utf-8")
    seed = int.from_bytes(hashlib.sha256(seed_src).digest()[:8], "big")
    rng = random.Random(seed)
    body = bytes(rng.randrange(256) for _ in range(720))
    # G6 — 페이로드 내 'PK' 시퀀스를 결정론 치환해 우연한 EOCD(PK\x05\x06) 형성을 차단.
    body = body.replace(b"PK", b"\x00K")
    path.write_bytes(b"PK\x03\x04" + body)
    # postcondition — 손상 파일은 반드시 유효 zip 이 아니어야 한다(verify ④축 손상 판정 전제).
    if zipfile.is_zipfile(path):
        raise bf.BFError(f"손상 파일이 유효 zip 으로 생성됨(EOCD 우연 포함): {doc['path']}")


def _render_lock(path: Path, doc: dict) -> None:
    """Office 잠금 임시파일 흉내(~$ 접두). 실무 공유폴더의 단골 쓰레기."""
    owner = doc.get("note", "unknown")
    path.write_bytes(b"\x08\x00" + owner.encode("utf-16-le") + b"\x00" * 148)


_RENDERERS = {
    "docx": _render_docx,
    "xlsx": _render_xlsx,
    "pptx": _render_pptx,
    "pdf": _render_pdf,
    "txt": _render_txt,
    "csv": _render_csv,
    "broken": _render_broken,
    "lock": _render_lock,
}


def generate(ledger_path: str, out_dir: str) -> list[str]:
    led = bf.load_ledger(ledger_path)
    out = Path(out_dir)
    _prepare_out(out)

    written: list[str] = []
    for doc in led["documents"]:
        rel = doc["path"]
        path = _abspath(out, rel)
        _RENDERERS[doc["type"]](path, doc)
        _set_mtime(path, doc["internal_date"])
        written.append(rel)
    return written


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="brain-fixture 데이터셋 생성기")
    ap.add_argument("ledger", help="원장 JSON 경로")
    ap.add_argument("--out", required=True, help="출력 폴더(비어있어야 함)")
    args = ap.parse_args(argv)
    try:
        written = generate(args.ledger, args.out)
    except bf.BFError as exc:
        print(f"[오류] {exc}", file=sys.stderr)
        return 2
    print(f"생성 완료: {len(written)}개 파일 → {args.out}")
    for rel in written:
        print(" -", rel)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
