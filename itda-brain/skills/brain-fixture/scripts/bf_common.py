#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""brain-fixture 공용 모듈 — 원장(ledger) 로드·검증, mtime, 문서 재파싱, 값 추출.

generate.py / verify.py / answer_sheet.py 가 공유한다. 원장이 SSoT 이므로
스키마 검증·값 추출 로직을 한 곳에 둔다(drift 차단).
"""
from __future__ import annotations

import json
import math
import os
import re
import sys
from datetime import datetime
from functools import reduce
from pathlib import Path

if sys.version_info[0] < 3:  # pragma: no cover
    sys.exit("Error: Python 3 필요. Windows 는 py -3 로 실행하세요.")

DOC_TYPES = {"docx", "xlsx", "pptx", "pdf", "txt", "csv", "broken", "lock"}
TRAP_TYPES = {
    "contradiction",
    "version-hell",
    "stale-rule",
    "time-warp",
    "decision-drift",
    "broken-file",
    "lock-file",
    "untitled",
}
BAIT_TYPES = {"direction", "scope", "duplicate"}
CONSISTENCY_OPS = {"sum", "product", "diff"}

# 인사이트(합성해야만 보이는 것) — REQ-050. type 은 자유 문자열(카탈로그는 문서화용, 검증 비강제).
INSIGHT_TYPE_CATALOG = {
    "negotiation-leverage",
    "threshold",
    "trend",
    "concentration",
    "margin",
    "deadline",
}
DERIV_NUMERIC_OPS = {"sum", "diff", "product"}  # 정수 결과
DERIV_RELATION_OPS = {"compare", "threshold"}  # 부등호 성립
DERIV_OPS = DERIV_NUMERIC_OPS | {"ratio"} | DERIV_RELATION_OPS
RELATIONS = {"lt", "gt", "lte", "gte", "eq", "ne"}


class BFError(Exception):
    """원장 스키마 위반 등 명시 에러 — 어느 필드가 왜 틀렸는지 담는다(no-silent-fallback)."""


# ---------------------------------------------------------------- 로드·검증

def _reject_json_constant(token: str):  # G4 — NaN/Infinity 명시 거부
    raise BFError(f"원장에 허용되지 않는 상수 '{token}' (NaN/Infinity 금지).")


def load_ledger(path: str | Path) -> dict:
    p = Path(path)
    if not p.is_file():
        raise BFError(f"원장 파일이 없습니다: {p}")
    try:
        data = json.loads(p.read_text(encoding="utf-8"), parse_constant=_reject_json_constant)
    except json.JSONDecodeError as exc:
        raise BFError(f"원장 JSON 파싱 실패({p}): {exc}") from exc
    if not isinstance(data, dict):
        raise BFError("원장 최상위는 객체(JSON object)여야 합니다.")
    validate_ledger(data)
    return data


def _need(obj: dict, key: str, where: str, types):
    if key not in obj:
        raise BFError(f"{where}: 필수 필드 '{key}' 누락")
    if not isinstance(obj[key], types):
        tn = types.__name__ if isinstance(types, type) else "/".join(t.__name__ for t in types)
        raise BFError(f"{where}: '{key}' 타입이 {tn} 이어야 합니다(현재 {type(obj[key]).__name__})")


def _finite_number(x) -> bool:
    # G4 — bool 은 수치가 아니며(True==1 회피), NaN/Infinity 도 배제.
    return isinstance(x, (int, float)) and not isinstance(x, bool) and math.isfinite(x)


def _need_number(obj: dict, key: str, where: str):
    if key not in obj:
        raise BFError(f"{where}: 필수 필드 '{key}' 누락")
    if not _finite_number(obj[key]):
        raise BFError(f"{where}: '{key}' 는 유한 숫자여야 합니다(bool·NaN·Infinity 불가, 현재 {obj[key]!r}).")


def _need_list(led: dict, key: str):
    # G5 — 키가 있으면 반드시 list. falsy({}·false·0·"") 를 SKIP 으로 오인하지 않는다.
    if key in led and not isinstance(led[key], list):
        raise BFError(f"ledger.{key}: 리스트여야 합니다(현재 {type(led[key]).__name__}). 미선언은 키 자체를 생략하세요.")


def normalized_relpath(path: str, where: str) -> str:
    """상대경로 안전성 검증 + 정규화 posix 경로 반환.

    절대경로·경로 탈출(``..``)을 명시 거부한다(출력 폴더 밖 파일 덮어쓰기 차단). 반환하는
    정규화 경로는 documents 유일성 판정에 쓴다(``a/../same.txt`` ↔ ``same.txt`` 물리 충돌 검출).
    ``~$`` 잠금 파일명 등 정상 상대경로(``..`` 없음)는 통과한다.
    """
    if not path or not isinstance(path, str):
        raise BFError(f"{where}: path 는 비어있지 않은 문자열이어야 합니다.")
    # 절대경로 거부 (POSIX '/', UNC '\\', Windows 드라이브 'C:').
    if path.startswith(("/", "\\")) or os.path.isabs(path) or (len(path) >= 2 and path[1] == ":"):
        raise BFError(f"{where}: 절대경로 금지 — 출력 폴더 기준 상대경로만 허용: '{path}'")
    norm = os.path.normpath(path.replace("\\", "/")).replace(os.sep, "/")
    if norm == ".." or norm.startswith("../") or norm.startswith("/"):
        raise BFError(f"{where}: 경로 탈출 금지 — 출력 폴더 밖으로 나가는 경로: '{path}'")
    return norm


def safe_join(root, rel: str) -> "Path":
    """rel 을 root 아래로 안전 결합한다. 절대경로·경로 탈출은 BFError.

    방어 심층 — 스키마 검증(normalized_relpath)과 별개로, 실제 파일 접근 직전 정규화 후
    루트 밖을 가리키면 차단한다(generate 쓰기·verify 읽기 공통 가드).
    """
    root_r = Path(root).resolve()
    target = (root_r / rel).resolve()
    try:
        target.relative_to(root_r)
    except ValueError as exc:
        raise BFError(f"경로 탈출: '{rel}' 가 대상 폴더 밖을 가리킵니다.") from exc
    return target


def validate_ledger(led: dict) -> None:
    """원장 전수 검증. 위반 시 BFError(어느 필드가 왜)."""
    _need(led, "profile", "ledger", dict)
    _need(led["profile"], "company", "profile", str)
    _need(led, "documents", "ledger", list)
    if not led["documents"]:
        raise BFError("ledger.documents: 최소 1개 문서가 필요합니다.")

    paths: set[str] = set()
    norm_seen: dict[str, str] = {}
    for i, doc in enumerate(led["documents"]):
        where = f"documents[{i}]"
        if not isinstance(doc, dict):
            raise BFError(f"{where}: 객체여야 합니다.")
        _need(doc, "path", where, str)
        _need(doc, "type", where, str)
        _need(doc, "internal_date", where, str)
        if doc["type"] not in DOC_TYPES:
            raise BFError(f"{where}: 알 수 없는 type '{doc['type']}' (허용: {sorted(DOC_TYPES)})")
        norm = normalized_relpath(doc["path"], where)  # 절대경로·경로 탈출 거부(FIX-1)
        if norm in norm_seen:
            raise BFError(
                f"{where}: 중복 경로 '{doc['path']}' — 정규화 시 '{norm_seen[norm]}' 와 같은 물리 파일(FIX-3)"
            )
        norm_seen[norm] = doc["path"]
        paths.add(doc["path"])
        _validate_internal_date(doc["internal_date"], where)
        _validate_doc_body(doc, where)

    # 선택 리스트 필드는 키가 있으면 반드시 list (G5 — falsy 를 SKIP 으로 오인 금지)
    for key in ("consistency", "traps", "baits", "insights"):
        _need_list(led, key)

    # consistency
    for i, c in enumerate(led.get("consistency", []) or []):
        where = f"consistency[{i}]"
        _need(c, "op", where, str)
        _need(c, "operands", where, list)
        _need_number(c, "expected", where)
        if c["op"] not in CONSISTENCY_OPS:
            raise BFError(f"{where}: 알 수 없는 op '{c['op']}' (허용: {sorted(CONSISTENCY_OPS)})")
        if not all(_finite_number(x) for x in c["operands"]):
            raise BFError(f"{where}: operands 는 유한 숫자 배열이어야 합니다(bool·NaN·Infinity 불가).")
        if not c["operands"]:
            raise BFError(f"{where}: operands 가 비어 있습니다.")
        if "expected_in" in c and c["expected_in"] not in paths:
            raise BFError(f"{where}: expected_in '{c['expected_in']}' 가 documents 에 없습니다.")

    # traps
    for i, t in enumerate(led.get("traps", []) or []):
        _validate_trap_like(t, f"traps[{i}]", paths, TRAP_TYPES)
    for i, b in enumerate(led.get("baits", []) or []):
        _validate_trap_like(b, f"baits[{i}]", paths, BAIT_TYPES)

    # insights (합성해야만 보이는 것) — REQ-050
    for i, ins in enumerate(led.get("insights", []) or []):
        _validate_insight(ins, f"insights[{i}]", paths)


def _validate_internal_date(s: str, where: str) -> None:
    try:
        datetime.strptime(s, "%Y-%m-%d %H:%M")
    except ValueError as exc:
        raise BFError(f"{where}: internal_date '{s}' 는 'YYYY-MM-DD HH:MM' 형식이어야 합니다.") from exc


def _validate_doc_body(doc: dict, where: str) -> None:
    t = doc["type"]
    if t == "docx":
        _need(doc, "title", where, str)
        _need(doc, "blocks", where, list)
        for j, blk in enumerate(doc["blocks"]):
            bw = f"{where}.blocks[{j}]"
            _need(blk, "kind", bw, str)
            if blk["kind"] in ("p", "h"):
                _need(blk, "text", bw, str)
            elif blk["kind"] == "table":
                _need(blk, "rows", bw, list)
            else:
                raise BFError(f"{bw}: kind 는 p|h|table 이어야 합니다(현재 '{blk['kind']}')")
    elif t == "xlsx":
        _need(doc, "sheets", where, dict)
        if not doc["sheets"]:
            raise BFError(f"{where}: sheets 가 비어 있습니다.")
    elif t == "pptx":
        _need(doc, "slides", where, list)
        if not doc["slides"]:
            raise BFError(f"{where}: slides 가 비어 있습니다.")
        for j, sl in enumerate(doc["slides"]):
            sw = f"{where}.slides[{j}]"
            _need(sl, "title", sw, str)
            _need(sl, "bullets", sw, list)
    elif t == "pdf":
        _need(doc, "title", where, str)
        _need(doc, "lines", where, list)
    elif t == "txt":
        _need(doc, "content", where, str)
    elif t == "csv":
        _need(doc, "rows", where, list)
    # broken / lock: 본문 없음(구조적 함정). note 는 선택.


def _validate_trap_like(t: dict, where: str, paths: set[str], allowed_types: set[str]) -> None:
    if not isinstance(t, dict):
        raise BFError(f"{where}: 객체여야 합니다.")
    _need(t, "id", where, str)
    _need(t, "type", where, str)
    _need(t, "title", where, str)
    _need(t, "detection", where, str)
    if t["type"] not in allowed_types:
        raise BFError(f"{where}: 알 수 없는 type '{t['type']}' (허용: {sorted(allowed_types)})")
    for tp in t.get("targets", []) or []:
        if tp not in paths:
            raise BFError(f"{where}: targets '{tp}' 가 documents 에 없습니다.")
    markers = t.get("markers", []) or []
    if not markers:
        # 유령 함정 차단(FIX-2) — 검증 가능한 marker 가 없으면 렌더되지 않은 함정이 정답지에 실린다.
        raise BFError(
            f"{where}: 검증 가능한 marker 가 최소 1개 필요합니다"
            " (path + value|text|unreadable|name_prefix)."
        )
    for j, m in enumerate(markers):
        mw = f"{where}.markers[{j}]"
        _need(m, "path", mw, str)
        if m["path"] not in paths:
            raise BFError(f"{mw}: path '{m['path']}' 가 documents 에 없습니다.")
        kinds = [k for k in ("value", "text", "unreadable", "name_prefix") if k in m]
        if len(kinds) != 1:
            raise BFError(f"{mw}: value|text|unreadable|name_prefix 중 정확히 1개를 지정해야 합니다(현재 {kinds}).")


def _validate_insight(ins: dict, where: str, paths: set[str]) -> None:
    if not isinstance(ins, dict):
        raise BFError(f"{where}: 객체여야 합니다.")
    _need(ins, "id", where, str)
    _need(ins, "type", where, str)  # 자유 문자열(카탈로그는 문서화용)
    _need(ins, "title", where, str)
    _need(ins, "conclusion", where, str)
    _need(ins, "surface_question", where, str)
    if "tier" not in ins:
        raise BFError(f"{where}: 필수 필드 'tier' 누락")
    if not (isinstance(ins["tier"], int) and not isinstance(ins["tier"], bool) and ins["tier"] in (1, 2, 3)):
        raise BFError(f"{where}: tier 는 1|2|3 정수여야 합니다(bool 불가, 현재 {ins['tier']!r}).")

    # evidence — 서로 다른 렌더 문서 ≥2
    _need(ins, "evidence", where, list)
    ev = ins["evidence"]
    if len(ev) < 2:
        raise BFError(f"{where}: evidence 는 서로 다른 문서 2개 이상 필요합니다(합성 강제).")
    if len(set(ev)) != len(ev):
        raise BFError(f"{where}: evidence 에 중복 문서가 있습니다 — 서로 다른 문서여야 합니다.")
    for e in ev:
        if not isinstance(e, str):
            raise BFError(f"{where}: evidence 원소는 문자열이어야 합니다.")
        if e not in paths:
            raise BFError(f"{where}: evidence '{e}' 가 documents 에 없습니다.")

    # derivation — 결론의 기계 도출식
    _need(ins, "derivation", where, dict)
    d = ins["derivation"]
    dw = f"{where}.derivation"
    _need(d, "op", dw, str)
    if d["op"] not in DERIV_OPS:
        raise BFError(f"{dw}: 알 수 없는 op '{d['op']}' (허용: {sorted(DERIV_OPS)})")
    _need(d, "operands", dw, list)
    ops = d["operands"]
    if not ops:
        raise BFError(f"{dw}: operands 가 비어 있습니다.")
    for j, o in enumerate(ops):
        ow = f"{dw}.operands[{j}]"
        _need_number(o, "value", ow)  # G4 — 유한 숫자만
        _need(o, "from", ow, str)
        if o["from"] not in ev:
            raise BFError(f"{ow}: from '{o['from']}' 가 evidence 에 없습니다(합성 근거는 evidence 문서여야 합니다).")
    # 피연산자가 서로 다른 evidence 문서 ≥2 에서 와야 진짜 '합성'
    if len({o["from"] for o in ops}) < 2:
        raise BFError(f"{dw}: 피연산자가 서로 다른 evidence 문서 2개 이상에서 와야 합니다(합성 강제).")

    op = d["op"]
    if op == "ratio":
        if len(ops) != 2:
            raise BFError(f"{dw}: ratio 는 피연산자가 정확히 2개여야 합니다(분자/분모).")
        if "scale" in d and not _finite_number(d["scale"]):
            raise BFError(f"{dw}: scale 은 유한 숫자여야 합니다.")
        if "round" in d and (not isinstance(d["round"], int) or isinstance(d["round"], bool)):
            raise BFError(f"{dw}: round 는 정수여야 합니다.")
        if ops[1]["value"] == 0:
            raise BFError(f"{dw}: ratio 분모가 0 입니다.")
    elif op in DERIV_RELATION_OPS:
        _need(d, "relation", dw, str)
        if d["relation"] not in RELATIONS:
            raise BFError(f"{dw}: 알 수 없는 relation '{d['relation']}' (허용: {sorted(RELATIONS)})")
        if len(ops) != 2:
            raise BFError(f"{dw}: {op} 는 피연산자가 정확히 2개여야 합니다.")

    # result — 기계 검증축(자연어 conclusion 과 분리). G2: 수치 결론을 relation op 으로 위장 차단.
    _need(ins, "result", where, dict)
    res = ins["result"]
    rw = f"{where}.result"
    _need(res, "kind", rw, str)
    if res["kind"] == "numeric":
        if op not in DERIV_NUMERIC_OPS and op != "ratio":
            raise BFError(f"{rw}: kind=numeric 은 수치 산출 op(sum·diff·product·ratio)여야 합니다(현재 op '{op}').")
        _need_number(res, "value", rw)  # derivation 재계산과 일치·스포일러 대상
    elif res["kind"] == "relation":
        if op not in DERIV_RELATION_OPS:
            raise BFError(f"{rw}: kind=relation 은 관계 op(compare·threshold)여야 합니다(현재 op '{op}').")
    else:
        raise BFError(f"{rw}: kind 는 numeric|relation 이어야 합니다(현재 '{res['kind']}').")


# ---------------------------------------------------------------- mtime

def mtime_ts(internal_date: str) -> float:
    """내부 날짜 문자열 → epoch(로컬 타임존). generate/verify 가 동일 산식을 쓴다."""
    return datetime.strptime(internal_date, "%Y-%m-%d %H:%M").timestamp()


# ---------------------------------------------------------------- 값 추출

_NUM_RE = re.compile(r"\d[\d,]*")


def extract_ints(text: str) -> set[int]:
    """텍스트에서 정수 토큰 집합을 뽑는다(콤마 제거). 'A4'→4 같은 잡음은 포함될 수 있으나
    membership(선언값 ∈ 집합) 검사만 쓰므로 안전(잡음은 false-absent 를 만들지 않음)."""
    out: set[int] = set()
    for tok in _NUM_RE.findall(text):
        tok = tok.strip(",")
        if not tok:
            continue
        try:
            out.add(int(tok.replace(",", "")))
        except ValueError:
            continue
    return out


def _cell_text(cell) -> str:
    if cell is None:
        return ""
    if isinstance(cell, float) and cell.is_integer():
        return str(int(cell))
    return str(cell)


def declared_text(doc: dict) -> str:
    """원장 문서 스펙에서 렌더될 전체 텍스트를 재구성한다(선언 기대값 추출용)."""
    t = doc["type"]
    parts: list[str] = []
    if t == "docx":
        parts.append(doc.get("title", ""))
        for blk in doc["blocks"]:
            if blk["kind"] in ("p", "h"):
                parts.append(blk["text"])
            elif blk["kind"] == "table":
                for row in blk["rows"]:
                    parts.extend(_cell_text(c) for c in row)
    elif t == "xlsx":
        for rows in doc["sheets"].values():
            for row in rows:
                parts.extend(_cell_text(c) for c in row)
    elif t == "pptx":
        for sl in doc["slides"]:
            parts.append(sl["title"])
            parts.extend(str(b) for b in sl["bullets"])
    elif t == "pdf":
        parts.append(doc.get("title", ""))
        parts.extend(str(x) for x in doc["lines"])
    elif t == "txt":
        parts.append(doc["content"])
    elif t == "csv":
        for row in doc["rows"]:
            parts.extend(_cell_text(c) for c in row)
    return "\n".join(parts)


def declared_ints(doc: dict) -> set[int]:
    return extract_ints(declared_text(doc))


# ---------------------------------------------------------------- 재파싱(독립)

class ParsedDoc:
    def __init__(self, text: str, ints: set[int], readable: bool, is_zip: bool):
        self.text = text
        self.ints = ints
        self.readable = readable  # 정상 문서로 열렸는가
        self.is_zip = is_zip


def reparse(path: Path, dtype: str) -> ParsedDoc:
    """생성된 파일을 원장과 독립적으로 재파싱한다. 실패 시 readable=False."""
    if not path.exists():
        return ParsedDoc("", set(), readable=False, is_zip=False)
    try:
        if dtype == "docx":
            return _reparse_docx(path)
        if dtype == "xlsx":
            return _reparse_xlsx(path)
        if dtype == "pptx":
            return _reparse_pptx(path)
        if dtype == "pdf":
            return _reparse_pdf(path)
        if dtype == "csv":
            return _reparse_csv(path)
        if dtype == "txt":
            txt = path.read_text(encoding="utf-8", errors="replace")
            return ParsedDoc(txt, extract_ints(txt), readable=True, is_zip=False)
        if dtype in ("broken", "lock"):
            import zipfile

            return ParsedDoc("", set(), readable=False, is_zip=zipfile.is_zipfile(path))
    except Exception:  # noqa: BLE001 — 재파싱 실패는 readable=False 로 표면화
        import zipfile

        return ParsedDoc("", set(), readable=False, is_zip=zipfile.is_zipfile(path))
    return ParsedDoc("", set(), readable=False, is_zip=False)


def _reparse_docx(path: Path) -> ParsedDoc:
    from docx import Document

    d = Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            parts.extend(cell.text for cell in row.cells)
    txt = "\n".join(parts)
    return ParsedDoc(txt, extract_ints(txt), readable=True, is_zip=True)


def _reparse_xlsx(path: Path) -> ParsedDoc:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=False, read_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            parts.extend(_cell_text(c) for c in row)
    wb.close()
    txt = "\n".join(parts)
    return ParsedDoc(txt, extract_ints(txt), readable=True, is_zip=True)


def _reparse_pptx(path: Path) -> ParsedDoc:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    txt = "\n".join(parts)
    return ParsedDoc(txt, extract_ints(txt), readable=True, is_zip=True)


def _reparse_csv(path: Path) -> ParsedDoc:
    """CSV 는 셀 단위로 파싱해 텍스트를 재구성한다. 구분자 콤마가 값의 천단위 콤마로
    오인돼 인접 필드가 병합되는 것을 막는다(예: '2026-01,100' → 병합 1100 방지)."""
    import csv as csvmod

    parts: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        for row in csvmod.reader(f):
            parts.extend(cell for cell in row)
    txt = "\n".join(parts)
    return ParsedDoc(txt, extract_ints(txt), readable=True, is_zip=False)


def _reparse_pdf(path: Path) -> ParsedDoc:
    import pypdf

    reader = pypdf.PdfReader(str(path))
    txt = "\n".join(page.extract_text() or "" for page in reader.pages)
    return ParsedDoc(txt, extract_ints(txt), readable=True, is_zip=True)


# ---------------------------------------------------------------- consistency 계산

def compute_consistency(op: str, operands: list) -> float:
    if op == "sum":
        return sum(operands)
    if op == "product":
        return reduce(lambda a, b: a * b, operands, 1)
    if op == "diff":
        return operands[0] - sum(operands[1:])
    raise BFError(f"알 수 없는 consistency op: {op}")


def as_int_if_integral(x):
    if isinstance(x, float) and x.is_integer():
        return int(x)
    return x


# ---------------------------------------------------------------- insight 도출식

def _relation_holds(rel: str, a, b) -> bool:
    return {
        "lt": a < b, "gt": a > b, "lte": a <= b, "gte": a >= b, "eq": a == b, "ne": a != b,
    }[rel]


def derivation_result(d: dict):
    """수치 산출 op(sum·diff·product·ratio)의 결과를 계산한다. ratio 는 scale·round 적용."""
    op = d["op"]
    vals = [o["value"] for o in d["operands"]]
    if op in ("sum", "diff", "product"):
        return compute_consistency(op, vals)
    if op == "ratio":
        scale = d.get("scale", 1)
        rnd = int(d.get("round", 0))
        return round(vals[0] / vals[1] * scale, rnd)
    raise BFError(f"derivation_result: 수치 산출 op 이 아닙니다 '{op}'")


def spoiler_signatures(d: dict, value) -> tuple[set[int], set[str]]:
    """스포일러 검사 대상 — 파생 결과가 어느 문서에도 직접 렌더되면 안 되는 (정수, 소수문자열).

    - 정수 채널: sum/diff/product 는 결과 정수. ratio 는 결과×10^round(예 13.7→137, 38.5→385) —
      소수 결론의 무점 인쇄("137")를 좁게 겨냥.
    - 소수 문자열 채널(G1): 소수 결론이 소수 표기("13.7"·"13.70")로 직접 렌더된 경우를 잡는다.
      재파싱 텍스트에 이 문자열이 있으면 FAIL. 포맷 변형(round~round+2 자리)을 포함한다.
    """
    op = d["op"]
    ints: set[int] = set()
    strs: set[str] = set()
    if op == "ratio":
        rnd = int(d.get("round", 0))
        ints.add(int(round(value * (10 ** rnd))))
        if not float(value).is_integer():  # 소수 결론 → 소수 문자열 표기도 금지
            for r in range(max(rnd, 1), max(rnd, 1) + 3):
                strs.add(f"{value:.{r}f}")
    else:
        ints.add(int(round(value)))
    return ints, strs
