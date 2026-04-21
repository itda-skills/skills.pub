"""Phase 3: python-pptx 기반 PPTX 패키징 모듈."""
from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

from itda_path import resolve_data_dir

# 16:9 표준 슬라이드 크기 (PowerPoint 와이드스크린 기본값)
_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)

# 슬라이드 파일명 패턴: slide_N.png
_SLIDE_PATTERN = "slide_*.png"
_SLIDE_NUM_RE = re.compile(r"slide_(\d+)")


def _get_blank_layout(prs: Presentation):
    """빈 슬라이드용 레이아웃을 반환한다.

    우선 이름이 "Blank"인 레이아웃을 찾고, 없으면 마지막 레이아웃으로 폴백한다.
    """
    for layout in prs.slide_layouts:
        if getattr(layout, "name", "").strip().lower() == "blank":
            return layout

    if len(prs.slide_layouts) == 0:
        raise ValueError("사용 가능한 슬라이드 레이아웃이 없습니다")

    return prs.slide_layouts[-1]


def build_pptx(
    image_paths: list[Path | str | None],
    output_path: Path | str | None = None,
) -> Path:
    """슬라이드 이미지 목록으로 PPTX를 생성한다.

    - 16:9 슬라이드 (13.333" × 7.5")
    - blank layout (slide_layouts[6])
    - 각 이미지를 Emu(0,0)에서 풀 블리드로 삽입
    - None 이미지는 빈 슬라이드로 추가
    - output_path 미지정 시 resolve_data_dir("slide-ai", "output") 사용

    Args:
        image_paths: 슬라이드로 삽입할 이미지 경로 목록. None이면 빈 슬라이드.
        output_path: 저장할 PPTX 경로. None이면 자동 생성.

    Returns:
        저장된 PPTX 파일의 Path.
    """
    prs = Presentation()
    prs.slide_width = _SLIDE_WIDTH
    prs.slide_height = _SLIDE_HEIGHT
    blank_layout = _get_blank_layout(prs)

    for image_path in image_paths:
        # blank layout으로 빈 슬라이드 추가
        slide = prs.slides.add_slide(blank_layout)
        if image_path is not None:
            slide.shapes.add_picture(
                str(image_path),
                left=Emu(0),
                top=Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

    # 출력 경로 결정
    if output_path is None:
        output_dir = resolve_data_dir("slide-ai", "output")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = output_dir / f"slides_{timestamp}.pptx"

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def rebuild_pptx(
    images_dir: Path | str,
    output_path: Path | str | None = None,
) -> Path:
    """기존 이미지 디렉토리에서 PPTX를 재빌드한다.

    images_dir에서 slide_*.png 파일을 숫자 순으로 정렬하여 PPTX 생성.

    Args:
        images_dir: slide_*.png 파일이 있는 디렉토리.
        output_path: 저장할 PPTX 경로. None이면 자동 생성.

    Returns:
        저장된 PPTX 파일의 Path.
    """
    images_dir = Path(images_dir)

    def _slide_number(p: Path) -> int:
        """slide_N.png에서 숫자 N을 추출한다. 없으면 0 반환."""
        m = _SLIDE_NUM_RE.search(p.stem)
        return int(m.group(1)) if m else 0

    # 숫자 순으로 정렬 (사전순이 아닌 숫자 기준)
    sorted_images = sorted(images_dir.glob(_SLIDE_PATTERN), key=_slide_number)

    return build_pptx(sorted_images, output_path=output_path)
