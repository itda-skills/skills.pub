#!/usr/bin/env python3
"""명단 파일 및 PPTX 템플릿 분석 스크립트.

# @MX:NOTE: JSON을 stdout 출력하여 Claude context 소모를 최소화함
# @MX:SPEC: SPEC-BADGE-001 FR-1, FR-2, FR-3, FR-4
"""

import argparse
import csv
import json
import math
import re
import sys
from pathlib import Path

if sys.version_info[0] < 3:
    sys.exit("Python 3 이상이 필요합니다.")

_PLACEHOLDER_RE = re.compile(r"\{[^}]+\}")
_INDEXED_RE = re.compile(r"^\{[^}_]+_(\d+)\}$")


def load_roster_xlsx(path):
    """openpyxl로 Excel(.xlsx) 파일을 로딩하여 (헤더 목록, 행 목록)을 반환."""
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    ws = wb.active
    if ws is None:
        wb.close()
        raise ValueError(f"활성 워크시트를 찾을 수 없습니다: {path}")
    rows = list(ws.iter_rows(values_only=True))
    wb.close()

    if not rows:
        raise ValueError(f"빈 명단 파일: {path}")

    headers = [str(h) if h is not None else "" for h in rows[0]]
    data = []
    for row in rows[1:]:
        record = {h: (str(v) if v is not None else "") for h, v in zip(headers, row)}
        data.append(record)

    if not data:
        raise ValueError(f"명단에 데이터 행이 없습니다: {path}")

    return headers, data


def load_roster_csv(path, encoding=None):
    """CSV 파일 로딩. UTF-8 → EUC-KR 순서로 인코딩을 시도.

    Args:
        path: CSV 파일 경로
        encoding: 강제 인코딩 지정 (None이면 자동 감지)
    """
    encodings = [encoding] if encoding else ["utf-8", "euc-kr"]

    for enc in encodings:
        try:
            with open(path, newline="", encoding=enc) as f:
                reader = csv.DictReader(f)
                rows = list(reader)
            if not rows:
                raise ValueError(f"빈 명단 파일: {path}")
            return list(rows[0].keys()), rows
        except UnicodeDecodeError:
            continue

    raise ValueError(f"CSV 인코딩 감지 실패 (시도: {encodings}): {path}")


def load_roster(path, encoding=None):
    """파일 확장자로 로더를 선택하여 명단을 로딩."""
    ext = Path(path).suffix.lower()
    if ext == ".xlsx":
        return load_roster_xlsx(path)
    elif ext == ".csv":
        return load_roster_csv(path, encoding=encoding)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext} (xlsx, csv만 지원)")


def extract_placeholders(prs):
    """모든 slide의 모든 shape에서 {xxx} 패턴 placeholder를 추출하여 정렬된 목록 반환."""
    found = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    found.update(_PLACEHOLDER_RE.findall(text))
    return sorted(found)


def detect_mode(placeholders):
    """placeholder 목록에서 모드(single/multi)와 슬라이드당 명찰 수를 반환.

    인덱스 방식({이름_1}, {이름_2})만 처리합니다.
    비인덱스 방식은 detect_badges_per_slide()를 사용하세요.

    Returns:
        (mode, badges_per_slide): ('single', 1) 또는 ('multi', N)
    """
    max_index = 0
    for ph in placeholders:
        m = _INDEXED_RE.match(ph)
        if m:
            max_index = max(max_index, int(m.group(1)))

    if max_index > 0:
        return "multi", max_index
    return "single", 1


def detect_badges_per_slide(slide):
    """슬라이드 객체에서 슬라이드당 명찰 수를 탐지.

    1순위: 인덱스 방식 ({이름_1}, {이름_2}) → max 인덱스
    2순위: 등장 횟수 방식 ({이름}이 4개 shape에 각각 존재) → max 등장 횟수

    Returns:
        (mode, badges_per_slide): ('single', 1) 또는 ('multi', N)
    """
    max_index = 0
    ph_shape_counts: dict[str, int] = {}

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # shape 내 고유 placeholder 집합 (같은 shape 내 중복은 1로 취급)
        shape_phs: set[str] = set()
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs)
            for ph in _PLACEHOLDER_RE.findall(text):
                shape_phs.add(ph)
                m = _INDEXED_RE.match(ph)
                if m:
                    max_index = max(max_index, int(m.group(1)))
        for ph in shape_phs:
            ph_shape_counts[ph] = ph_shape_counts.get(ph, 0) + 1

    if max_index > 0:
        return "multi", max_index

    if ph_shape_counts:
        badges = max(ph_shape_counts.values())
        if badges > 1:
            return "multi", badges

    return "single", 1


def suggest_mapping(placeholders, columns):
    """placeholder의 base name과 컬럼명을 비교하여 자동 매핑을 제안.

    Returns:
        dict: {base_name: column_name} — 중괄호/인덱스 제거된 base name이 키
    """
    mapping = {}
    seen = set()
    for ph in placeholders:
        base = re.sub(r"_\d+$", "", ph.strip("{}"))
        if base not in seen and base in columns:
            mapping[base] = base
            seen.add(base)
    return mapping


def analyze(data_path, template_path, encoding=None):
    """명단 파일과 PPTX 템플릿을 분석하여 결과 dict를 반환.

    Args:
        data_path: 명단 파일 경로 (.xlsx 또는 .csv)
        template_path: PPTX 템플릿 파일 경로
        encoding: CSV 인코딩 강제 지정 (None이면 자동 감지)

    Raises:
        ValueError: 파일 없음, 빈 파일, 슬라이드 수 오류 등
    """
    from pptx import Presentation

    columns, rows = load_roster(data_path, encoding=encoding)

    prs = Presentation(template_path)
    slide_count = len(prs.slides)
    if slide_count != 1:
        raise ValueError(
            f"템플릿은 슬라이드 1개여야 합니다 (현재: {slide_count}개)"
        )

    placeholders = extract_placeholders(prs)
    mode, badges_per_slide = detect_badges_per_slide(prs.slides[0])
    mapping = suggest_mapping(placeholders, columns)

    total_people = len(rows)
    total_slides = math.ceil(total_people / badges_per_slide)
    last_slide_people = total_people % badges_per_slide
    last_slide_empty = (badges_per_slide - last_slide_people) % badges_per_slide

    return {
        "roster": {
            "file": Path(data_path).name,
            "format": Path(data_path).suffix.lower().lstrip("."),
            "columns": columns,
            "row_count": total_people,
            "sample": rows[:3],
        },
        "template": {
            "file": Path(template_path).name,
            "slide_count": slide_count,
            "placeholders": placeholders,
            "badges_per_slide": badges_per_slide,
            "mode": mode,
        },
        "mapping_suggestion": mapping,
        "output_estimate": {
            "total_people": total_people,
            "badges_per_slide": badges_per_slide,
            "total_slides": total_slides,
            "last_slide_empty_slots": last_slide_empty,
        },
    }


def main():
    parser = argparse.ArgumentParser(description="명단 파일 및 PPTX 템플릿 분석")
    parser.add_argument(
        "--data", required=True, help="명단 파일 경로 (.xlsx 또는 .csv)"
    )
    parser.add_argument("--template", required=True, help="PPTX 템플릿 파일 경로")
    parser.add_argument(
        "--encoding", help="CSV 인코딩 강제 지정 (예: utf-8, euc-kr)"
    )
    args = parser.parse_args()

    try:
        result = analyze(args.data, args.template, encoding=args.encoding)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except Exception as e:
        print(json.dumps({"error": str(e)}, ensure_ascii=False), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
