#!/usr/bin/env python3
"""명찰 PPTX 생성 스크립트.

# @MX:NOTE: XML deep copy로 슬라이드를 복제하여 원본 서식(폰트/크기/색상/위치) 완전 보존
# @MX:SPEC: SPEC-BADGE-001 FR-5, FR-6
"""

import argparse
import copy
import csv
import json
import re
import sys
from datetime import datetime
from pathlib import Path

if sys.version_info[0] < 3:
    sys.exit("Python 3 이상이 필요합니다.")

# 기본 컬럼 매핑: placeholder base name → Excel/CSV 컬럼명
DEFAULT_MAPPING = {"이름": "이름", "부서": "부서", "회사명": "회사명"}

_INDEXED_RE = re.compile(r"^\{[^}_]+_(\d+)\}$")


# ── 명단 로딩 (analyze.py와 동일한 로직) ──────────────────────────────────────


def load_roster_xlsx(path):
    """openpyxl로 Excel(.xlsx) 파일 로딩."""
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    ws = wb.active
    if ws is None:
        wb.close()
        raise ValueError(f"활성 워크시트를 찾을 수 없습니다: {path}")
    rows = list(ws.iter_rows(values_only=True))
    wb.close()

    if not rows:
        raise ValueError(f"빈 명단 파일: {path}")

    headers = [str(h) if h is not None else "" for h in rows[0]]
    data = []
    for row in rows[1:]:
        record = {h: (str(v) if v is not None else "") for h, v in zip(headers, row)}
        data.append(record)

    if not data:
        raise ValueError(f"명단에 데이터 행이 없습니다: {path}")

    return headers, data


def load_roster_csv(path, encoding=None):
    """CSV 파일 로딩. UTF-8 → EUC-KR 순서로 인코딩 시도."""
    encodings = [encoding] if encoding else ["utf-8", "euc-kr"]

    for enc in encodings:
        try:
            with open(path, newline="", encoding=enc) as f:
                reader = csv.DictReader(f)
                rows = list(reader)
            if not rows:
                raise ValueError(f"빈 명단 파일: {path}")
            return list(rows[0].keys()), rows
        except UnicodeDecodeError:
            continue

    raise ValueError(f"CSV 인코딩 감지 실패 (시도: {encodings}): {path}")


def load_roster(path, encoding=None):
    """파일 확장자로 로더를 선택하여 명단 로딩."""
    ext = Path(path).suffix.lower()
    if ext == ".xlsx":
        return load_roster_xlsx(path)
    elif ext == ".csv":
        return load_roster_csv(path, encoding=encoding)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext} (xlsx, csv만 지원)")


# ── 슬라이드 복제 ──────────────────────────────────────────────────────────────


def duplicate_slide(prs, template_slide):
    """템플릿 슬라이드를 XML deep copy로 복제하여 새 슬라이드를 추가.

    # @MX:NOTE: add_slide 후 기본 요소를 제거하고 원본 XML을 복사하는 방식으로 서식 완전 보존
    """
    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # 새 슬라이드의 기본 요소 제거
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # 원본 슬라이드 XML 복사
    for shape in template_slide.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shape._element))

    return new_slide


# ── 텍스트 치환 ────────────────────────────────────────────────────────────────


def replace_text_in_shape(shape, replacements):
    """shape 내 각 paragraph의 run을 합쳐 치환하고 첫 번째 run에 기록.

    paragraph 단위로 run 텍스트를 합산하여 placeholder가 여러 run에
    분할된 경우도 올바르게 처리한다. 치환 후 첫 번째 run에 결과를,
    나머지 run은 빈 문자열로 설정하여 서식 손실을 최소화한다.
    """
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.runs:
            continue
        full_text = "".join(run.text for run in paragraph.runs)
        new_text = full_text
        for old, new_val in replacements.items():
            new_text = new_text.replace(old, new_val)
        if new_text != full_text:
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""


# ── 치환 dict 생성 ─────────────────────────────────────────────────────────────


def build_replacements(chunk, mapping, mode, badges_per_slide):
    """청크 데이터와 매핑으로 {placeholder: 값} 치환 dict를 생성.

    Args:
        chunk: 이번 슬라이드에 들어갈 행 목록
        mapping: {base_name: column_name} 매핑
        mode: 'single' 또는 'multi'
        badges_per_slide: 슬라이드당 명찰 수
    """
    replacements = {}
    if mode == "single":
        row = chunk[0] if chunk else {}
        for ph_base, col in mapping.items():
            replacements["{" + ph_base + "}"] = str(row.get(col, ""))
    else:  # multi
        for i in range(1, badges_per_slide + 1):
            row = chunk[i - 1] if i <= len(chunk) else {}
            for ph_base, col in mapping.items():
                key = "{" + ph_base + "_" + str(i) + "}"
                replacements[key] = str(row.get(col, ""))
    return replacements


# ── 비인덱스 다중 명찰: 위치 기반 badge 그룹핑 ────────────────────────────────


def _compute_badge_groups(template_slide, badges_per_slide):
    """비인덱스 다중 명찰 슬라이드에서 badge 단위 shape 인덱스 그룹을 반환.

    # @MX:NOTE: anchor placeholder(가장 많이 등장하는 것)를 기준으로
    #           각 shape를 가장 가까운 anchor에 할당하는 nearest-anchor 클러스터링.
    #           인덱스 없는 {이름}×4 형태의 비인덱스 템플릿 전용.

    Returns:
        list[list[int]]: badges_per_slide 길이의 shape 인덱스 그룹 목록.
                         그룹 순서: anchor 위치 (top, left) 오름차순.
    """
    ph_re = re.compile(r"\{[^}]+\}")
    shapes_list = list(template_slide.shapes)

    ph_shape_entries: list[tuple[int, object]] = []  # (index, shape)
    ph_counts: dict[str, int] = {}

    for idx, shape in enumerate(shapes_list):
        if not shape.has_text_frame:
            continue
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            continue
        text = "".join("".join(r.text for r in p.runs) for p in tf.paragraphs)
        phs = set(ph_re.findall(text))
        if phs:
            ph_shape_entries.append((idx, shape))
            for ph in phs:
                ph_counts[ph] = ph_counts.get(ph, 0) + 1

    if not ph_shape_entries:
        return [[] for _ in range(badges_per_slide)]

    # anchor: 가장 많이 등장하는 placeholder (= badges_per_slide번 등장해야 함)
    anchor_ph = max(ph_counts, key=lambda k: ph_counts[k])

    anchor_entries: list[tuple[int, object]] = [
        (idx, shape) for idx, shape in ph_shape_entries
        if anchor_ph in "".join(
            "".join(r.text for r in p.runs)
            for p in shape.text_frame.paragraphs
        )
    ]
    anchor_entries.sort(key=lambda x: (x[1].top, x[1].left))

    if len(anchor_entries) != badges_per_slide:
        # anchor 개수 불일치 시 빈 그룹 반환 (호출자가 fallback 처리)
        return []

    # 각 ph_shape를 nearest anchor에 할당
    groups: list[list[int]] = [[] for _ in range(badges_per_slide)]
    for shape_idx, shape in ph_shape_entries:
        cx = shape.left + shape.width / 2
        cy = shape.top + shape.height / 2

        min_dist = float("inf")
        closest = 0
        for group_i, (_, anchor) in enumerate(anchor_entries):
            ax = anchor.left + anchor.width / 2
            ay = anchor.top + anchor.height / 2
            dist = (cx - ax) ** 2 + (cy - ay) ** 2
            if dist < min_dist:
                min_dist = dist
                closest = group_i
        groups[closest].append(shape_idx)

    return groups


# ── 슬라이드 삭제 ─────────────────────────────────────────────────────────────


def _remove_first_slide(prs):
    """프레젠테이션의 첫 번째 슬라이드(템플릿 원본)를 제거."""
    from pptx.oxml.ns import qn

    xml_slides = prs.slides._sldIdLst
    if not len(xml_slides):
        return
    first_elem = xml_slides[0]
    rId = first_elem.get(qn("r:id"))
    if rId:
        prs.part.drop_rel(rId)
    first_elem.getparent().remove(first_elem)


# ── 메인 생성 함수 ─────────────────────────────────────────────────────────────


def generate(data_path, template_path, output_path, mapping, encoding=None, badges_per_slide=None):
    """명찰 PPTX를 생성하여 output_path에 저장.

    # @MX:ANCHOR: 스킬의 핵심 생성 함수. 매핑, 모드, 슬라이드 복제 모두 여기서 조율
    # @MX:REASON: 외부에서 직접 호출되는 메인 API (CLI main + SKILL.md 워크플로 양쪽)

    Args:
        data_path: 명단 파일 경로
        template_path: PPTX 템플릿 파일 경로
        output_path: 출력 PPTX 파일 경로
        mapping: {base_placeholder_name: column_name} 매핑 dict
        encoding: CSV 인코딩 강제 지정
        badges_per_slide: 슬라이드당 명찰 수. None이면 템플릿에서 자동 감지.

    Returns:
        dict: output, total_slides, total_people, badges_per_slide

    Raises:
        ValueError: 빈 명단, 슬라이드 수 오류 등
    """
    from pptx import Presentation

    _, rows = load_roster(data_path, encoding=encoding)
    if not rows:
        raise ValueError("명단이 비어 있습니다.")

    prs = Presentation(template_path)
    if len(prs.slides) != 1:
        raise ValueError(
            f"템플릿은 슬라이드 1개여야 합니다 (현재: {len(prs.slides)}개)"
        )

    template_slide = prs.slides[0]

    if badges_per_slide is None:
        # analyze.py를 거치지 않은 직접 호출 시 폴백: 템플릿에서 자동 감지
        placeholder_re = re.compile(r"\{[^}]+\}")
        ph_shape_counts: dict[str, int] = {}
        max_index = 0
        for shape in template_slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            shape_phs: set[str] = set()
            for para in tf.paragraphs:
                text = "".join(run.text for run in para.runs)
                for ph in placeholder_re.findall(text):
                    shape_phs.add(ph)
                    m = _INDEXED_RE.match(ph)
                    if m:
                        max_index = max(max_index, int(m.group(1)))
            for ph in shape_phs:
                ph_shape_counts[ph] = ph_shape_counts.get(ph, 0) + 1

        if max_index > 0:
            badges_per_slide = max_index
        elif ph_shape_counts:
            badges_per_slide = max(ph_shape_counts.values())
        else:
            badges_per_slide = 1

    mode = "multi" if badges_per_slide > 1 else "single"

    # 비인덱스 다중 명찰 여부 감지: 템플릿 슬라이드에서 직접 확인
    _ph_re_check = re.compile(r"\{[^}]+\}")
    _all_phs: set[str] = set()
    for _shape in template_slide.shapes:
        if _shape.has_text_frame:
            for _para in _shape.text_frame.paragraphs:
                _all_phs.update(_ph_re_check.findall("".join(r.text for r in _para.runs)))
    has_indexed = any(_INDEXED_RE.match(ph) for ph in _all_phs)

    # 매핑 키가 템플릿에 존재하는지 검증
    missing = []
    for ph_base in mapping:
        if has_indexed:
            if not any(ph.startswith("{" + ph_base + "_") for ph in _all_phs):
                missing.append("{" + ph_base + "_N}")
        else:
            if "{" + ph_base + "}" not in _all_phs:
                missing.append("{" + ph_base + "}")
    if missing:
        raise ValueError(
            f"템플릿에서 placeholder를 찾을 수 없습니다: {', '.join(missing)}"
        )

    badge_groups: list[list[int]] = []
    if mode == "multi" and not has_indexed:
        badge_groups = _compute_badge_groups(template_slide, badges_per_slide)

    total_people = len(rows)
    chunks = [
        rows[i : i + badges_per_slide] for i in range(0, total_people, badges_per_slide)
    ]

    for chunk in chunks:
        new_slide = duplicate_slide(prs, template_slide)

        if badge_groups:
            # 비인덱스 다중 명찰: shape 그룹별로 각 사람 데이터를 치환
            new_shapes = list(new_slide.shapes)
            for person_idx, row in enumerate(chunk):
                if person_idx >= len(badge_groups):
                    break
                replacements = {
                    "{" + ph_base + "}": str(row.get(col, ""))
                    for ph_base, col in mapping.items()
                }
                for shape_idx in badge_groups[person_idx]:
                    replace_text_in_shape(new_shapes[shape_idx], replacements)
        else:
            # 단일 또는 인덱스 다중: 키 기반 치환
            replacements = build_replacements(chunk, mapping, mode, badges_per_slide)
            for shape in new_slide.shapes:
                replace_text_in_shape(shape, replacements)

    # 템플릿 원본(첫 번째 슬라이드) 제거
    _remove_first_slide(prs)

    prs.save(output_path)

    return {
        "output": str(output_path),
        "total_slides": len(chunks),
        "total_people": total_people,
        "badges_per_slide": badges_per_slide,
    }


def main():
    parser = argparse.ArgumentParser(description="명찰 PPTX 생성")
    parser.add_argument("--data", required=True, help="명단 파일 경로")
    parser.add_argument("--template", required=True, help="PPTX 템플릿 파일 경로")
    parser.add_argument("--output", help="출력 파일 경로 (기본: 현재 경로)")
    parser.add_argument(
        "--badges-per-slide",
        type=int,
        help="슬라이드당 명찰 수 (analyze.py의 badges_per_slide 값). 미지정 시 템플릿에서 자동 감지.",
    )
    parser.add_argument(
        "--mapping",
        help=(
            '컬럼 매핑 JSON (예: \'{"이름":"성명","회사명":"기관"}\'). '
            "미지정 시 기본값 사용: 이름→이름, 부서→부서, 회사명→회사명. "
            "지정 시 기본값에 merge됩니다."
        ),
    )
    parser.add_argument(
        "--encoding", help="CSV 인코딩 강제 지정 (예: utf-8, euc-kr)"
    )
    args = parser.parse_args()

    mapping = dict(DEFAULT_MAPPING)  # 기본값 복사
    if args.mapping:
        try:
            override = json.loads(args.mapping)
        except json.JSONDecodeError as e:
            print(
                json.dumps({"error": f"매핑 JSON 파싱 오류: {e}"}, ensure_ascii=False),
                file=sys.stderr,
            )
            sys.exit(1)
        mapping.update(override)  # 기본값에 사용자 지정 merge

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    if args.output:
        p = Path(args.output)
        output_path = str(p.parent / f"{p.stem}_{timestamp}{p.suffix}")
    else:
        output_path = str(Path.cwd() / f"itda-name-badge-{timestamp}.pptx")

    try:
        result = generate(
            args.data, args.template, output_path, mapping,
            encoding=args.encoding,
            badges_per_slide=args.badges_per_slide,
        )
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except Exception as e:
        print(json.dumps({"error": str(e)}, ensure_ascii=False), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
