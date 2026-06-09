"""deckkit — pptx-design 스킬의 공개 헬퍼 API (SPEC-PPTX-DESIGN-001 REQ-006 / -002 REQ-001·003·006).

per-invocation 생성 스크립트(gen.py)가 import 해서 쓰는 안정 API.
python-pptx + matplotlib + Pillow 기반. 크로스플랫폼(macOS/Linux), Office 불필요.

핵심 설계:
- CJK 안전 폰트: set_run_font 가 latin/ea/cs typeface 를 동시 지정(한글 tofu/세리프 폴백 회피).
- ★한글 타이포 가드(REQ-001): 한글 run 에 음수/과대 자간·비-한글(세리프/thin 라틴) 폰트가 지정되면
  자동으로 안전 고딕·0 자간으로 클램프. force=True 로 우회 가능.
- ★LibreOffice 안정 폰트 체인(REQ-006): 실측상 또렷한 굵은 고딕(Noto Sans KR·Pretendard)을 우선
  선택. NanumGothic·Apple SD Gothic Neo 는 LibreOffice 에서 자간 벌어진 명조풍으로 치환되어 후순위.
- 네이티브 편집 차트(REQ-003): add_native_chart 가 python-pptx add_chart 로 PowerPoint 편집가능
  차트 객체를 생성하고 DESIGN.md 팔레트·CJK 안전 폰트를 적용. matplotlib 래스터 차트는 고디자인 옵션.
- 그라디언트/모티프: python-pptx 는 그라디언트 fill 미지원 → Pillow PNG 로 베이크해 풀블리드 임베드.
- 도형 기본 파랑 테두리/그림자: rect() 가 기본 제거.
"""
from __future__ import annotations
import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

EMU_PER_IN = 914400
DEFAULT_W_IN = 13.333   # 16:9 와이드
DEFAULT_H_IN = 7.5

# 한글 run 자간 상한(pt). 음수는 0 으로, 이 값 초과는 이 값으로 클램프(REQ-001).
KR_SPACING_CAP_PT = 2.0

# LibreOffice 치환 안전 한글 폰트 후보 — ★렌더 또렷함(굵은 고딕) 우선 정렬(SPEC-PPTX-DESIGN-002 §1 실측).
#   Noto Sans KR·Pretendard = 또렷한 굵은 고딕 / NanumGothic·Apple SD Gothic Neo = 자간 벌어진 명조풍 치환.
# (경로 템플릿, LibreOffice 패밀리명). 존재하는 첫 항목을 kr_font_name() 이 반환.
_KR_FONT_FILES = [
    ("/Users/{user}/Library/Fonts/NotoSansKR[wght].ttf", "Noto Sans KR"),
    ("/Users/{user}/Library/Fonts/NotoSansKR-VariableFont_wght.ttf", "Noto Sans KR"),
    ("/usr/share/fonts/opentype/noto/NotoSansKR-Regular.otf", "Noto Sans KR"),
    ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", "Noto Sans CJK KR"),
    ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", "Noto Sans CJK KR"),
    ("/Users/{user}/Library/Fonts/Pretendard-Regular.otf", "Pretendard"),
    ("/usr/share/fonts/truetype/pretendard/Pretendard-Regular.otf", "Pretendard"),
    ("/System/Library/Fonts/AppleSDGothicNeo.ttc", "Apple SD Gothic Neo"),
    ("C:/Windows/Fonts/malgun.ttf", "Malgun Gothic"),
    ("/Users/{user}/Library/Fonts/NanumGothic-Regular.ttf", "NanumGothic"),
    ("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", "NanumGothic"),
]

# matplotlib 차트 라벨용 폰트 — matplotlib 은 LibreOffice 치환 버그가 없어 어떤 고딕이든 또렷.
# 변수 폰트(Noto Sans KR[wght]) 보다 정적 파일이 reliable → 정적 우선 정렬.
_MPL_FONT_FILES = [
    ("/Users/{user}/Library/Fonts/Pretendard-Regular.otf", "Pretendard"),
    ("/Users/{user}/Library/Fonts/NanumGothic-Regular.ttf", "NanumGothic"),
    ("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", "NanumGothic"),
    ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", "Noto Sans CJK KR"),
    ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", "Noto Sans CJK KR"),
    ("/System/Library/Fonts/AppleSDGothicNeo.ttc", "Apple SD Gothic Neo"),
    ("/Users/{user}/Library/Fonts/NotoSansKR[wght].ttf", "Noto Sans KR"),
]

# 한글 run 에 지정해도 안전한(=Korean-capable) 폰트 allowlist(소문자 비교).
# 이 목록 밖 폰트(라틴 전용 세리프/디스플레이/산세)가 한글 run 에 오면 kr_font_name() 으로 강제.
_KR_CAPABLE_FONTS = {
    "noto sans kr", "noto sans cjk kr", "noto sans cjk", "noto sans cjk jp",
    "noto serif kr", "noto serif cjk kr",
    "pretendard", "pretendard variable", "pretendard jp",
    "apple sd gothic neo", "applegothic", "apple gothic", "applemyungjo",
    "nanumgothic", "nanum gothic", "nanumbarungothic", "nanum barun gothic",
    "nanumsquare", "nanumsquare neo", "nanummyeongjo", "nanum myeongjo",
    "malgun gothic", "malgungothic", "gulim", "굴림", "batang", "바탕",
    "dotum", "돋움", "gungsuh", "spoqa han sans neo", "spoqa han sans",
    "ibm plex sans kr", "gothic a1", "gmarket sans", "gmarketsans",
    "source han sans", "source han sans kr", "source han sans k",
    "함초롬바탕", "함초롬돋움", "맑은 고딕", "나눔고딕",
}


def hexrgb(h: str) -> RGBColor:
    """'cc785c' 또는 '#cc785c' → RGBColor."""
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def new_deck(w_in: float = DEFAULT_W_IN, h_in: float = DEFAULT_H_IN) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    return prs


def blank_slide(prs: Presentation):
    """완전 빈 레이아웃(idx 6) 슬라이드 추가."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, hexcolor: str):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = hexrgb(hexcolor)


def _kill_shadow(sp):
    """LibreOffice 가 preset effectRef 로 강제 상속하는 드롭섀도를 XML 로 중화."""
    spPr = sp._element.spPr
    for tag in ("a:effectLst", "a:effectDag"):
        existing = spPr.find(qn(tag))
        if existing is not None:
            spPr.remove(existing)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))
    try:
        sp.shadow.inherit = False
    except Exception:
        pass


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=False):
    """도형 추가. 기본: 테두리 없음·그림자 없음. fill/line 은 hex 문자열.

    radius: ROUNDED_RECTANGLE 의 모서리 비율(0~0.5). shadow=True 면 기본 그림자 유지.
    """
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = hexrgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = hexrgb(line); sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = float(radius)
        except Exception:
            pass
    if not shadow:
        _kill_shadow(sp)
    return sp


def has_hangul(s) -> bool:
    """문자열에 한글(완성형·자모)이 하나라도 있으면 True."""
    return bool(re.search(r"[가-힣ㄱ-ㅎㅏ-ㅣ]", s or ""))


# 하위호환 alias(내부/verify 용)
_has_hangul = has_hangul


def is_kr_capable_font(name) -> bool:
    """폰트명이 한글을 안전하게 렌더하는(Korean-capable) allowlist 에 있으면 True."""
    return bool(name) and name.strip().lower() in _KR_CAPABLE_FONTS


def set_run_font(run, name=None, size=None, bold=None, italic=None, color=None,
                 spacing=None, force=False):
    """run 폰트 설정. ★name 지정 시 latin/ea/cs 동시 적용(한글이 해당 폰트로 렌더).

    ★한글 타이포 가드(SPEC-PPTX-DESIGN-002 REQ-001) — force=False(기본):
      - 한글 포함 run 에 비-한글 폰트(라틴 세리프/thin 디스플레이 등)를 지정하면
        kr_font_name() 의 안전 고딕으로 강제(세리프/명조/붓글씨 폴백 차단).
      - 한글 run 의 음수 자간(spacing<0)은 0 으로, KR_SPACING_CAP_PT 초과는 캡으로 클램프
        (LibreOffice 에서 음수 자간이 오히려 벌어지는 결함 차단).
      - 라틴 전용 run 은 DESIGN.md 의 음수 자간·세리프 디스플레이를 그대로 허용.
    force=True 면 가드를 우회한다(고급 사용자가 의도적으로 적용할 때만).

    주의: 가드가 run.text 로 한글을 판별하므로 set_run_font 는 run.text 설정 후 호출한다(현 toolkit 규약).
    """
    is_kr = has_hangul(run.text)
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = hexrgb(color)
    if name is not None:
        if is_kr and not force and not is_kr_capable_font(name):
            name = kr_font_name()   # 한글에 비안전 폰트 → 안전 고딕으로 강제
        f.name = name
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {}); rPr.append(el)
            el.set("typeface", name)
    if spacing is not None:
        if is_kr and not force:
            spacing = max(0.0, min(float(spacing), KR_SPACING_CAP_PT))  # 음수→0, 과대→캡
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, line_spacing=None, space_after=None):
    """runs: [(text, {font kwargs}), ...]. 단일 문단. 정렬 어긋남 방지 위해 margin=0."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if space_after is not None:
        p.space_after = Pt(space_after)
    for text, kw in runs:
        r = p.add_run(); r.text = text; set_run_font(r, **kw)
    return tb


def add_paragraph(text_frame, runs, align=PP_ALIGN.LEFT, space_after=None, line_spacing=None):
    """기존 text_frame 에 문단 추가(멀티라인 본문용)."""
    p = text_frame.add_paragraph()
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    for text, kw in runs:
        r = p.add_run(); r.text = text; set_run_font(r, **kw)
    return p


def add_table(slide, x, y, w, h, nrows, ncols, col_w=None):
    """빈 표 추가 후 table 객체 반환. 셀 스타일은 호출측에서 cell.fill/run 으로."""
    gt = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Inches(cw)
    gt.first_row = False
    gt.horz_banding = False
    return gt


def add_picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w is not None:
        kw["width"] = Inches(w)
    if h is not None:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)


# ---------- Pillow: 그라디언트/모티프 ----------
def linear_gradient(w, h, c_top, c_bottom, vertical=True):
    """선형 그라디언트 PIL.Image. c_*: (r,g,b). python-pptx 그라디언트 미지원 대응."""
    import numpy as np
    from PIL import Image
    t = np.array(c_top, float); b = np.array(c_bottom, float)
    n = h if vertical else w
    ramp = (t[None, :] + (b - t)[None, :] * np.linspace(0, 1, n)[:, None]).astype("uint8")
    arr = (np.repeat(ramp[:, None, :], w, axis=1) if vertical
           else np.repeat(ramp[None, :, :], h, axis=0))
    return Image.fromarray(arr, "RGB")


def radial_glow(w, h, base, glow, cx=0.5, cy=0.4, radius=0.7):
    """방사형 글로우 PIL.Image(어두운 base 위 glow 색). c: (r,g,b)."""
    import numpy as np
    from PIL import Image
    yy, xx = np.mgrid[0:h, 0:w].astype(float)
    d = np.sqrt(((xx / w - cx)) ** 2 + ((yy / h - cy)) ** 2) / radius
    a = np.clip(1 - d, 0, 1)[:, :, None]
    base = np.array(base, float); glow = np.array(glow, float)
    arr = (base[None, None, :] * (1 - a) + glow[None, None, :] * a).astype("uint8")
    return Image.fromarray(arr, "RGB")


# ---------- matplotlib 한글 ----------
def mpl_korean():
    """matplotlib 한글 폰트 등록 + 마이너스 정상화. 차트 그리기 전 1회 호출.

    matplotlib 은 LibreOffice 치환 버그가 없어 정적 고딕(Pretendard/Nanum/Noto)이면 또렷.
    반환: 등록된 폰트 패밀리명(없으면 'DejaVu Sans')."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager as fm
    user = os.environ.get("USER", "")
    chosen = None
    for tmpl, fam in _MPL_FONT_FILES:
        path = tmpl.format(user=user)
        if os.path.exists(path):
            try:
                fm.fontManager.addfont(path)
                chosen = fam
                break
            except Exception:
                continue
    plt.rcParams["font.family"] = chosen or "DejaVu Sans"
    plt.rcParams["axes.unicode_minus"] = False
    return chosen


def kr_font_name():
    """이 환경에서 LibreOffice 가 또렷한 굵은 고딕으로 렌더하는 한글 폰트명(set_run_font name 용).

    REQ-006 결정론: _KR_FONT_FILES 순서(Noto Sans KR·Pretendard 우선)대로 존재하는 첫 폰트를 선택.
    어떤 후보도 없으면 'Noto Sans KR'(LibreOffice 가 시스템 CJK 로 폴백)을 반환."""
    user = os.environ.get("USER", "")
    for tmpl, fam in _KR_FONT_FILES:
        if os.path.exists(tmpl.format(user=user)):
            return fam
    return "Noto Sans KR"


# ---------- 네이티브 편집 차트 (REQ-003) ----------
_CHART_KINDS = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "stacked_column_100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "stacked_bar_100": XL_CHART_TYPE.BAR_STACKED_100,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "radar": XL_CHART_TYPE.RADAR,
}
_LINE_KINDS = {"line", "line_markers", "area", "area_stacked", "radar"}
_CIRCULAR_KINDS = {"pie", "doughnut"}
_LEGEND_POS = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "top": XL_LEGEND_POSITION.TOP,
    "right": XL_LEGEND_POSITION.RIGHT,
    "left": XL_LEGEND_POSITION.LEFT,
    "corner": XL_LEGEND_POSITION.CORNER,
}


def _chart_cjk_font(font_obj, name, size=None, color=None, bold=None):
    """차트 텍스트(범례·축·라벨) 폰트에 CJK 안전 typeface 를 latin/ea/cs 동시 지정."""
    if name:
        font_obj.name = name
    if size is not None:
        font_obj.size = Pt(size)
    if bold is not None:
        font_obj.bold = bold
    if color is not None:
        font_obj.color.rgb = hexrgb(color)
    if name:
        rPr = font_obj._rPr
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {}); rPr.append(el)
            el.set("typeface", name)


def add_native_chart(slide, x, y, w, h, kind, categories, series,
                     palette=None, font_name=None, font_size=11, font_color=None,
                     legend=True, legend_pos="bottom", gridlines=False,
                     data_labels=False, number_format=None, chart_title=None,
                     value_axis=True, category_axis=True,
                     gap_width=None, overlap=None, line_width=2.5):
    """네이티브 편집가능 pptx 차트 추가(python-pptx add_chart) — PowerPoint 에서 차트 객체로 편집 가능.

    matplotlib 래스터 차트(이미지 임베드, 편집 불가)와 달리 수치가 그대로 남아 받는 사람이
    PowerPoint 에서 데이터를 직접 수정할 수 있다. DESIGN.md 팔레트·CJK 안전 폰트가 적용된다.

    인자:
      kind: 'column'|'bar'|'line'|'line_markers'|'pie'|'doughnut'|'stacked_column'|
            'stacked_bar'|'stacked_column_100'|'area'|'area_stacked'|'radar' (또는 XL_CHART_TYPE).
      categories: x축 카테고리 라벨 리스트.
      series: [(name, [values...]), ...]. pie/doughnut 은 단일 series 의 point 별 색.
      palette: series(또는 원형 차트의 point) 별 hex 색 리스트(기본 matplotlib 파랑/주황 추방).
      font_name: 범례·축·라벨 폰트(None → kr_font_name() 안전 고딕, latin/ea/cs 동시 지정).
      data_labels/number_format/legend/gridlines/value_axis/category_axis: 표시 옵션.
    반환: graphic_frame(추가 조작은 gf.chart 로).
    """
    ct = _CHART_KINDS.get(kind) if isinstance(kind, str) else kind
    if ct is None:
        raise ValueError(f"알 수 없는 chart kind: {kind!r}. 지원: {sorted(_CHART_KINDS)}")
    is_circular = (kind in _CIRCULAR_KINDS) if isinstance(kind, str) else False
    is_line = (kind in _LINE_KINDS) if isinstance(kind, str) else False

    cd = CategoryChartData()
    cd.categories = list(categories)
    for nm, vals in series:
        cd.add_series(nm, tuple(vals), number_format=number_format)
    gf = slide.shapes.add_chart(ct, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart

    fam = font_name or kr_font_name()
    _chart_cjk_font(chart.font, fam, size=font_size, color=font_color)

    # 차트 제목 — 기본 끔(덱이 별도 제목 텍스트를 그리는 경우가 많아 중복 회피). 주면 표시.
    chart.has_title = bool(chart_title)
    if chart_title:
        chart.chart_title.text_frame.text = chart_title
        _chart_cjk_font(chart.chart_title.text_frame.paragraphs[0].font, fam,
                        size=(font_size or 11) + 2, color=font_color)

    chart.has_legend = bool(legend)
    if legend:
        chart.legend.position = _LEGEND_POS.get(legend_pos, XL_LEGEND_POSITION.BOTTOM)
        chart.legend.include_in_layout = False
        _chart_cjk_font(chart.legend.font, fam, size=font_size, color=font_color)

    plot = chart.plots[0]

    # 팔레트 적용(기본 파랑/주황 추방)
    if palette:
        if is_circular:
            plot.vary_by_categories = True
            pts = plot.series[0].points
            for i, pt in enumerate(pts):
                col = palette[i % len(palette)]
                pt.format.fill.solid(); pt.format.fill.fore_color.rgb = hexrgb(col)
        else:
            for i, ser in enumerate(chart.series):
                col = palette[i % len(palette)]
                if is_line:
                    ser.format.line.color.rgb = hexrgb(col)
                    ser.format.line.width = Pt(line_width)
                else:
                    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = hexrgb(col)
                    ser.format.line.fill.background()   # 막대 테두리 제거(AI 슬라이드 티 회피)

    if gap_width is not None:
        try: plot.gap_width = gap_width
        except Exception: pass
    if overlap is not None:
        try: plot.overlap = overlap
        except Exception: pass

    if data_labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        if number_format:
            dl.number_format = number_format
            dl.number_format_is_linked = False
        _chart_cjk_font(dl.font, fam, size=font_size, color=font_color)

    # 축(원형 차트는 축이 없어 예외 무시)
    try:
        ca = chart.category_axis
        ca.has_major_gridlines = False
        ca.has_minor_gridlines = False
        if not category_axis:
            ca.visible = False
        _chart_cjk_font(ca.tick_labels.font, fam, size=font_size, color=font_color)
    except Exception:
        pass
    try:
        va = chart.value_axis
        va.has_major_gridlines = bool(gridlines)
        va.has_minor_gridlines = False
        if not value_axis:
            va.visible = False
        _chart_cjk_font(va.tick_labels.font, fam, size=font_size, color=font_color)
    except Exception:
        pass

    return gf


def save_deck(prs, path):
    os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
    prs.save(path)
    return path
