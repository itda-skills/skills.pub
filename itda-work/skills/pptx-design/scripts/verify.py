"""verify — PPTX 배치/렌더 3층 + 타이포 자동 검증 (SPEC-PPTX-DESIGN-001 REQ-005 / -002 REQ-005).

층:
  (A) 지오메트리(python-pptx): 경계이탈(off-slide)·퇴화도형 [HARD] · 텍스트박스 겹침 [advisory]
  (B) 콘텐츠 대조: --tokens 파일의 필수 토큰이 pptx 텍스트에 존재 [HARD]
  (C) OCR 렌더 대조(tesseract kor+eng, 다크 반전): 산출률·한글명 [advisory, 형식오탐 회피용]
  (D) ★한글 타이포 정적 검사: 한글 run 의 음수/과대 자간·비안전(세리프/라틴 디스플레이) 폰트
      [advisory] — 미적 결함(자간 벌어짐·명조 폴백)의 자동 트립와이어(SPEC-PPTX-DESIGN-002 REQ-005).
  (E) ★스타일 휴리스틱: AI 기본값 안티패턴 — 구분 부호(·/—/–) 남발·좌측 액센트 바 남발·
      수직 중앙 몰림 [advisory] — "엉성한 스타일" 자동 트립와이어(SPEC-PPTX-DESIGN-003 REQ-105).
  (F) ★wrap 유발 오버플로(python-pptx): 긴 문단이 좁은 박스에서 래핑돼 필요 높이가 박스를 넘는 경우
      [advisory] — view_issues L1(문단수×폰트×1.2)의 사각지대를 폭 근사로 보완(#413).

HARD GATE = (경계이탈 + 퇴화도형 + 빈슬라이드 + 토큰누락) == 0  → PASS 시 exit 0. (타이포·스타일·wrap 은 advisory)

사용:
  python3 verify.py <pptx> [--tokens tokens.txt] [--ko 삼성전자,하이닉스]
                          [--out DIR] [--dpi 110] [--no-ocr]
"""
import os
import re
import sys
import json
import glob
import argparse
import math

from pptx import Presentation
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import render as render_mod  # noqa: E402
import deckkit as dk  # noqa: E402  (has_hangul · is_kr_capable_font · KR_SPACING_CAP_PT 재사용)

EMU = 914400
TOL = int(0.04 * EMU)
YIELD_MIN_CHARS = 90
YIELD_RATIO = 0.22
OCR_CFG = "--oem 1 --psm 11 -l kor+eng"


def _ocr_available():
    try:
        import pytesseract  # noqa
        import shutil
        return shutil.which("tesseract") is not None
    except Exception:
        return False


def bbox(sp):
    if None in (sp.left, sp.top, sp.width, sp.height):
        return None
    return (int(sp.left), int(sp.top), int(sp.left) + int(sp.width), int(sp.top) + int(sp.height))


def area(b):
    return max(0, b[2] - b[0]) * max(0, b[3] - b[1])


def inter(a, b):
    x0, y0 = max(a[0], b[0]), max(a[1], b[1])
    x1, y1 = min(a[2], b[2]), min(a[3], b[3])
    return max(0, x1 - x0) * max(0, y1 - y0)


def shape_text(sp):
    try:
        if sp.has_text_frame:
            return sp.text_frame.text
    except Exception:
        pass
    return ""


def hangul_alnum(s):
    return re.sub(r"[^0-9A-Za-z가-힣]", "", s)


def slide_texts(prs):
    per = {}
    for si, sl in enumerate(prs.slides, start=1):
        buf = []
        for sp in sl.shapes:
            try:
                if sp.has_text_frame:
                    buf.append(sp.text_frame.text)
            except Exception:
                pass
            try:
                if sp.has_table:
                    for r in sp.table.rows:
                        for c in r.cells:
                            buf.append(c.text)
            except Exception:
                pass
        per[si] = "\n".join(buf)
    return per


def run_typo_issues(prs):
    """(D) 한글 타이포 정적 검사 — 한글 run 의 음수/과대 자간·비안전 폰트 advisory 수집.

    deckkit 가드(set_run_font)는 생성 시 이를 교정하지만, 이 검사는
      (1) force=True 우회·수기 빌드 덱의 회귀,
      (2) 가드 없이 만든 구버전 결함 덱
    을 적발한다. 모두 advisory(HARD GATE 비차단).
    """
    neg, wide, unsafe = [], [], []
    for si, sl in enumerate(prs.slides, start=1):
        for sp in sl.shapes:
            try:
                if not sp.has_text_frame:
                    continue
            except Exception:
                continue
            for para in sp.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text or ""
                    if not dk.has_hangul(t):
                        continue
                    rPr = run._r.find(qn("a:rPr"))
                    if rPr is None:
                        continue
                    spc_raw = rPr.get("spc")
                    if spc_raw is not None:
                        try:
                            spc_pt = int(spc_raw) / 100.0
                        except ValueError:
                            spc_pt = None
                        if spc_pt is not None and spc_pt < 0:
                            neg.append({"slide": si, "text": t[:24].replace("\n", " "),
                                        "spc_pt": round(spc_pt, 2)})
                        elif spc_pt is not None and spc_pt > dk.KR_SPACING_CAP_PT:
                            wide.append({"slide": si, "text": t[:24].replace("\n", " "),
                                         "spc_pt": round(spc_pt, 2)})
                    face = None
                    for tag in ("a:ea", "a:latin", "a:cs"):
                        el = rPr.find(qn(tag))
                        if el is not None and el.get("typeface"):
                            face = el.get("typeface")
                            break
                    if face is not None and not dk.is_kr_capable_font(face):
                        unsafe.append({"slide": si, "text": t[:24].replace("\n", " "), "font": face})
    return neg, wide, unsafe


# ── (E) 스타일 휴리스틱 임계 (SPEC-PPTX-DESIGN-003 REQ-105, 튜닝 가능) ──
STYLE_PUNCT_CHARS = "·—–"          # 구분 부호: 가운데점·em대시·en대시
STYLE_PUNCT_MIN_PT = 18.0          # 이 크기 이상(또는 크기 미지정) run 만 카운트 — 메타/캡션 줄 제외 proxy
STYLE_PUNCT_MAX_PER_SLIDE = 5      # 슬라이드당 초과 시 advisory
STYLE_BAR_MAX_W_IN = 0.12          # 좌측 액센트 바 후보: 폭 ≤
STYLE_BAR_MIN_H_IN = 0.35          #                       높이 ≥
STYLE_BAR_SNAP_IN = 0.08           # 다른 도형 왼쪽 모서리 흡착 판정 거리
STYLE_BAR_MAX_PER_DECK = 3         # 덱 전체 초과 시에만 advisory(절제 사용 1~3곳은 무죄)
STYLE_VC_MAX_H_RATIO = 0.55        # 콘텐츠 union 높이 < 55%H 이면서
STYLE_VC_TOL_RATIO = 0.08          # union 중심이 캔버스 중심 ±8%H → 중앙 몰림
STYLE_VC_MIN_SHAPES = 3


def run_style_issues(prs):
    """(E) 스타일 휴리스틱 — 디자인 시스템 없는 AI 기본값 안티패턴 advisory 3종.

    ① style_punct_overuse: 큰 글씨 run 의 '·'/'—'/'–' 합이 슬라이드당 임계 초과(부호 남발).
    ② style_edge_bar_overuse: 다른 도형 왼쪽 모서리에 흡착된 얇은 세로 바가 덱 임계 초과
       (좌측 액센트 라인 남발 — "적당해야 세련". 임계 이하 절제 사용은 보고하지 않음).
    ③ style_v_center_cram: 비배경 도형 union 이 낮은 높이로 캔버스 수직 중앙에 몰림.
    모두 advisory — HARD GATE 산식 불변. 임계는 모듈 상수로 노출(오탐 시 튜닝).
    """
    W, H = int(prs.slide_width), int(prs.slide_height)
    SA = W * H
    punct, bars, cram = [], [], []
    for si, sl in enumerate(prs.slides, start=1):
        boxes = []      # 비배경 도형 bbox
        bar_cands = []  # 좌측 바 후보 bbox
        pcount = 0
        for sp in sl.shapes:
            b = bbox(sp)
            if b is None:
                continue
            w_in, h_in = (b[2] - b[0]) / EMU, (b[3] - b[1]) / EMU
            if w_in <= 0 or h_in <= 0:
                continue
            is_bg = area(b) >= 0.92 * SA and b[0] <= TOL and b[1] <= TOL
            if not is_bg:
                boxes.append(b)
                if w_in <= STYLE_BAR_MAX_W_IN and h_in >= STYLE_BAR_MIN_H_IN:
                    bar_cands.append(b)
            try:
                has_tf = sp.has_text_frame
            except Exception:
                has_tf = False
            if has_tf:
                for para in sp.text_frame.paragraphs:
                    for run in para.runs:
                        sz = run.font.size
                        if sz is not None and sz.pt < STYLE_PUNCT_MIN_PT:
                            continue
                        pcount += sum((run.text or "").count(ch) for ch in STYLE_PUNCT_CHARS)
        if pcount > STYLE_PUNCT_MAX_PER_SLIDE:
            punct.append({"slide": si, "count": pcount})
        snap = int(STYLE_BAR_SNAP_IN * EMU)
        for bb in bar_cands:
            bw = bb[2] - bb[0]
            for ob in boxes:
                if ob == bb:
                    continue
                if (ob[2] - ob[0]) <= bw * 2:
                    continue
                if abs(ob[0] - bb[2]) <= snap:
                    ov = min(bb[3], ob[3]) - max(bb[1], ob[1])
                    if ov >= 0.5 * (bb[3] - bb[1]):
                        bars.append({"slide": si})
                        break
        if len(boxes) >= STYLE_VC_MIN_SHAPES:
            u0 = min(b[1] for b in boxes)
            u1 = max(b[3] for b in boxes)
            uh = u1 - u0
            ucy = (u0 + u1) / 2
            if uh < STYLE_VC_MAX_H_RATIO * H and abs(ucy - H / 2) < STYLE_VC_TOL_RATIO * H:
                cram.append({"slide": si, "h_ratio": round(uh / H, 2),
                             "center_off_ratio": round(abs(ucy - H / 2) / H, 3)})
    if len(bars) <= STYLE_BAR_MAX_PER_DECK:
        bars = []
    return punct, bars, cram


# ── (F) wrap 유발 오버플로 임계 (#413, advisory — 폰트 메트릭 근사라 보수적) ──
WRAP_LINE_HEIGHT = 1.2        # 기본 줄높이 계수 (view_issues L1 과 동일)
WRAP_OVERFLOW_TOL = 1.05      # 추정 필요높이가 박스높이의 이 배수 초과 시 신호(em 근사 오차 여유)
WRAP_DEFAULT_PT = 18.0        # run.font.size 미상속 시 보수적 기본 글자크기
WRAP_MIN_BOX_W_IN = 0.5       # 장식용 초협소 박스 제외
PT_EMU = 12700               # 1pt = 12700 EMU (72pt = 1in = 914400 EMU)


def _est_line_em(s):
    """문자 클래스 기반 1줄 텍스트 폭 근사(단위 em = font_size 배). 폰트 비의존·결정론.

    정확한 글리프 폭 대신 CJK=1.0 / 라틴 대문자·숫자=0.62 / 소문자=0.5 / 공백·기호=0.3~0.42
    의 평균 전진폭으로 근사한다. advisory 목적(픽셀 정밀이 아니라 wrap 위험 포착)에 충분하며,
    시스템 폰트 로드(비결정)를 피해 동일 입력 → 동일 출력을 보장한다.
    """
    w = 0.0
    for ch in s:
        o = ord(ch)
        if ch == "\t":
            w += 2.0
        elif (0xAC00 <= o <= 0xD7A3) or (0x3000 <= o <= 0x30FF) or \
             (0x4E00 <= o <= 0x9FFF) or (0xFF00 <= o <= 0xFFEF):
            w += 1.0                      # 한글·CJK·가나·전각
        elif ch == " ":
            w += 0.30
        elif ch.isupper() or ch.isdigit():
            w += 0.62
        elif ch.islower():
            w += 0.50
        else:
            w += 0.42                     # 라틴 구두점·기타
    return w


def run_wrap_issues(prs):
    """(F) wrap 유발 오버플로 — 폭 보정 줄 수로 필요 높이를 재추정해 박스 초과를 적발.

    view_issues 의 L1 휴리스틱(needed = 문단수 × 폰트 × 1.2)은 긴 단일 문단이 좁은 박스에서
    여러 줄로 래핑돼도 1줄로 계산해 wrap 유발 오버플로를 놓친다(#413). 여기서는 각 문단의
    추정 폭(_est_line_em)을 박스 가용 폭으로 나눠 줄 수를 보정한 뒤 필요 높이를 다시 잰다.

    신호 조건: (보정 필요높이 > 가용높이 × TOL) AND (보정 줄 수 합 > 문단 수)
      - 후자가 "naive(문단수)로는 통과하나 wrap 으로 넘침"이라는 #413 사각지대를 정확히 가둔다.
    advisory — 폭 근사 특성상 보수적 임계. HARD GATE 산식 불변(오탐이 덱 빌드를 깨지 않음).
    word_wrap=False(넘침은 OOB 가 담당)·AutoSize=TextToFitShape(글자 축소로 맞춤)·배경 도형은 제외.
    SHAPE_TO_FIT_TEXT(python-pptx textbox 기본값)는 박스가 세로로 늘어나 선언 높이를 넘겨 아래
    요소를 침범하므로(절대좌표라 자동으로 안 밀림) 검사 대상에 포함한다.
    """
    from pptx.enum.text import MSO_AUTO_SIZE
    out = []
    W, H = int(prs.slide_width), int(prs.slide_height)
    SA = W * H
    min_w = int(WRAP_MIN_BOX_W_IN * EMU)
    for si, sl in enumerate(prs.slides, start=1):
        for sp in sl.shapes:
            try:
                if not sp.has_text_frame:
                    continue
            except Exception:
                continue
            b = bbox(sp)
            if b is None:
                continue
            box_w, box_h = b[2] - b[0], b[3] - b[1]
            if box_w < min_w or box_h <= 0:
                continue
            if area(b) >= 0.92 * SA and b[0] <= TOL and b[1] <= TOL:
                continue                  # 배경 도형
            tf = sp.text_frame
            try:
                if tf.word_wrap is False:
                    continue
            except Exception:
                pass
            try:
                if tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
                    continue          # 글자 크기를 박스에 맞춰 축소 → 오버플로 없음
            except Exception:
                pass
            ml = int(tf.margin_left) if tf.margin_left is not None else int(0.1 * EMU)
            mr = int(tf.margin_right) if tf.margin_right is not None else int(0.1 * EMU)
            mt = int(tf.margin_top) if tf.margin_top is not None else int(0.05 * EMU)
            mb = int(tf.margin_bottom) if tf.margin_bottom is not None else int(0.05 * EMU)
            usable_w, usable_h = box_w - ml - mr, box_h - mt - mb
            if usable_w <= 0 or usable_h <= 0:
                continue
            cap_pt = usable_w / PT_EMU    # 가용 폭(pt)
            naive_lines = wrapped_lines = 0
            needed_pt = 0.0
            for para in tf.paragraphs:
                ptxt = "".join((r.text or "") for r in para.runs) or (para.text or "")
                psz = 0.0
                for r in para.runs:
                    if r.font.size is not None:
                        psz = max(psz, r.font.size.pt)
                if psz <= 0:
                    psz = WRAP_DEFAULT_PT
                naive_lines += 1
                if not ptxt.strip():
                    wrapped_lines += 1
                    needed_pt += psz * WRAP_LINE_HEIGHT
                    continue
                plines = max(1, math.ceil((_est_line_em(ptxt) * psz) / cap_pt)) if cap_pt > 0 else 1
                wrapped_lines += plines
                needed_pt += plines * psz * WRAP_LINE_HEIGHT
            needed_emu = needed_pt * PT_EMU
            if needed_emu > usable_h * WRAP_OVERFLOW_TOL and wrapped_lines > naive_lines:
                out.append({"slide": si, "text": shape_text(sp)[:40].replace("\n", " "),
                            "naive_lines": naive_lines, "est_lines": wrapped_lines,
                            "box_h_in": round(box_h / EMU, 2), "needed_in": round(needed_emu / EMU, 2)})
    return out


def ocr_blob(im):
    import numpy as np
    import pytesseract
    from PIL import ImageOps
    g = im.convert("L")
    g = g.resize((g.width * 2, g.height * 2))
    mean = float(np.asarray(g).mean())
    txt = pytesseract.image_to_string(g, config=OCR_CFG)
    if mean < 130:
        txt += "\n" + pytesseract.image_to_string(ImageOps.invert(g), config=OCR_CFG)
    return txt


def verify(pptx_path, tokens=None, ko=None, out_dir=None, dpi=110, do_ocr=True):
    prs = Presentation(pptx_path)
    W, H = int(prs.slide_width), int(prs.slide_height)
    SA = W * H
    stem = os.path.splitext(os.path.basename(pptx_path))[0]
    out_dir = out_dir or os.path.join(os.path.dirname(os.path.abspath(pptx_path)), "_verify")
    os.makedirs(out_dir, exist_ok=True)

    issues = {"out_of_bounds": [], "zero_size": [], "blank_slide": [], "missing_content": [],
              "text_overlap": [], "ocr_low_yield": [], "ocr_missing_korean": [],
              "kr_neg_spacing": [], "kr_wide_spacing": [], "kr_unsafe_font": [],
              "style_punct_overuse": [], "style_edge_bar_overuse": [], "style_v_center_cram": [],
              "text_wrap_overflow": [], "render_unavailable": []}
    perslide = {}
    sltext = slide_texts(prs)

    # (D) 한글 타이포 정적 검사 (advisory)
    neg, wide, unsafe = run_typo_issues(prs)
    issues["kr_neg_spacing"] = neg
    issues["kr_wide_spacing"] = wide
    issues["kr_unsafe_font"] = unsafe

    # (E) 스타일 휴리스틱 (advisory) — AI 기본값 안티패턴 트립와이어
    s_punct, s_bars, s_cram = run_style_issues(prs)
    issues["style_punct_overuse"] = s_punct
    issues["style_edge_bar_overuse"] = s_bars
    issues["style_v_center_cram"] = s_cram

    # (F) wrap 유발 오버플로 (advisory) — view_issues L1 사각지대 보완(#413)
    issues["text_wrap_overflow"] = run_wrap_issues(prs)

    # (A) 지오메트리
    for si, sl in enumerate(prs.slides, start=1):
        flags = []
        texts = []
        for sp in sl.shapes:
            b = bbox(sp)
            if b is None:
                continue
            w, h = b[2] - b[0], b[3] - b[1]
            a = area(b)
            is_bg = a >= 0.92 * SA and b[0] <= TOL and b[1] <= TOL
            if (w < 0 or h < 0) or (w == 0 and h == 0):
                issues["zero_size"].append({"slide": si}); flags.append((b, "zero")); continue
            if w <= 0 or h <= 0:
                continue
            over = []
            if b[0] < -TOL: over.append(f"left {(-b[0])/EMU:.2f}in")
            if b[1] < -TOL: over.append(f"top {(-b[1])/EMU:.2f}in")
            if b[2] > W + TOL: over.append(f"right +{(b[2]-W)/EMU:.2f}in")
            if b[3] > H + TOL: over.append(f"bottom +{(b[3]-H)/EMU:.2f}in")
            if over and not is_bg:
                issues["out_of_bounds"].append({"slide": si, "over": over,
                                                "text": shape_text(sp)[:40].replace("\n", " ")})
                flags.append((b, "oob"))
            if is_bg:
                continue
            tx = shape_text(sp).strip()
            if tx:
                texts.append((b, len(tx)))
        for i in range(len(texts)):
            for j in range(i + 1, len(texts)):
                bi, bj = texts[i][0], texts[j][0]
                if min(texts[i][1], texts[j][1]) < 3:
                    continue
                ov = inter(bi, bj)
                if ov <= 0:
                    continue
                ai, aj = area(bi), area(bj)
                if max(ai, aj) / max(1, min(ai, aj)) >= 6:
                    continue
                if ov / min(ai, aj) > 0.70:
                    issues["text_overlap"].append({"slide": si, "ratio": round(ov / min(ai, aj), 2)})
                    flags.append((bi, "ov")); flags.append((bj, "ov"))
        perslide[si] = flags

    # (B) 콘텐츠 대조
    if tokens:
        alltext = "\n".join(sltext.values()).replace(" ", "")
        for tok in tokens:
            if tok.replace(" ", "") not in alltext:
                issues["missing_content"].append(tok)

    # 렌더 + (C) OCR + 빈슬라이드 + 주석
    try:
        jpgs = render_mod.render(pptx_path, out_dir=os.path.join(out_dir, f"{stem}_render"), dpi=dpi)
    except Exception as e:
        # 렌더 도구(soffice/pdftoppm) 부재·실패는 인프라 문제이지 덱 결함이 아니다(#621).
        # blank_slide(HARD GATE 항목)가 아니라 비차단 advisory(render_unavailable)로 분리한다.
        # 렌더 의존 검사(OCR 층 C + 이미지 기반 빈슬라이드 탐지)는 jpgs 가 비어 자연히 생략된다.
        jpgs = []
        issues["render_unavailable"].append(
            {"note": f"렌더 도구 사용 불가 — 렌더 의존 검사(OCR·이미지 빈슬라이드) 생략: {e}"})
    if jpgs and len(jpgs) < len(list(prs.slides)):
        issues["blank_slide"].append({"note": f"렌더 {len(jpgs)}/{len(list(prs.slides))}장"})

    ocr_on = do_ocr and _ocr_available() and bool(jpgs)
    deck_raw = ""
    annot = []
    ocrflag = {}
    for si, jp in enumerate(jpgs, start=1):
        im = Image.open(jp).convert("RGB")
        small = im.resize((80, 45))
        cols = small.getcolors(80 * 45) or []
        if cols and max(cols, key=lambda x: x[0])[0] / (80 * 45) > 0.992:
            issues["blank_slide"].append({"slide": si})
        if ocr_on:
            raw = ocr_blob(im)
            deck_raw += hangul_alnum(raw)
            ptx = len(hangul_alnum(sltext.get(si, "")))
            ocr = len(hangul_alnum(raw))
            if ptx >= YIELD_MIN_CHARS and ocr / max(1, ptx) < YIELD_RATIO:
                issues["ocr_low_yield"].append({"slide": si, "pptx_chars": ptx, "ocr_chars": ocr,
                                                "ratio": round(ocr / max(1, ptx), 2)})
                ocrflag[si] = round(ocr / max(1, ptx), 2)
        annot.append((si, im))
    if ocr_on and ko:
        for k in ko:
            if k.replace(" ", "") not in deck_raw:
                issues["ocr_missing_korean"].append(k)

    # 주석 몽타주
    if annot:
        tw, COLS = 600, 3
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial Bold.ttf", 24)
        except Exception:
            font = ImageFont.load_default()
        thumbs = [(si, im.resize((tw, int(tw * im.height / im.width)))) for si, im in annot]
        th = thumbs[0][1].height
        rows = (len(thumbs) + COLS - 1) // COLS
        pad, lab = 8, 30
        cv = Image.new("RGB", (COLS * tw + (COLS + 1) * pad, rows * (th + lab) + (rows + 1) * pad), (235, 235, 235))
        dd = ImageDraw.Draw(cv)
        sx, sy = tw / W, th / H
        for k, (si, t) in enumerate(thumbs):
            c, r = k % COLS, k // COLS
            x, y = pad + c * (tw + pad), pad + r * (th + lab + pad)
            for (b, tag) in perslide.get(si, []):
                col = {"oob": (230, 30, 30), "ov": (255, 140, 0), "zero": (200, 0, 200)}.get(tag, (230, 30, 30))
                dd.rectangle([x + b[0] * sx, y + lab + b[1] * sy, x + b[2] * sx, y + lab + b[3] * sy], outline=col, width=4)
            cv.paste(t, (x, y + lab))
            g = len(perslide.get(si, []))
            o = ocrflag.get(si)
            head = f"S{si}" + (f"  G{g}" if g else "") + (f"  OCR{o}" if o is not None else "")
            dd.rectangle([x, y, x + tw, y + lab - 3], fill=(180, 0, 0) if (g or o) else (20, 20, 20))
            dd.text((x + 6, y + 3), head, fill=(255, 255, 255), font=font)
        cv.save(os.path.join(out_dir, f"{stem}-annotated.png"))

    counts = {k: len(v) for k, v in issues.items()}
    gate = counts["out_of_bounds"] + counts["zero_size"] + counts["blank_slide"] + counts["missing_content"]
    result = {"pptx": pptx_path, "hard_gate_pass": gate == 0, "counts": counts,
              "ocr_ran": ocr_on, "issues": issues}
    # ensure_ascii=False 로 한국어/특수문자(em-dash 등)를 그대로 쓰므로 utf-8 고정.
    # (Windows 기본 cp949 로 열면 비-cp949 문자에서 UnicodeEncodeError, #621)
    json.dump(result, open(os.path.join(out_dir, f"{stem}.json"), "w", encoding="utf-8"),
              ensure_ascii=False, indent=2)
    return result


def main():
    # Windows 콘솔(cp949 등)에서 한국어·em-dash(—) 출력이 UnicodeEncodeError 로 죽지 않도록
    # stdout/stderr 를 utf-8 로 강제한다(#621). 리다이렉트/비-TTY 등으로 실패하면 무시.
    for _stream in (sys.stdout, sys.stderr):
        try:
            _stream.reconfigure(encoding="utf-8")
        except Exception:
            pass
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--tokens", help="필수 토큰 파일(1줄 1토큰). 미지정 시 콘텐츠 대조 생략")
    ap.add_argument("--ko", help="OCR 한글명 advisory(쉼표 구분)")
    ap.add_argument("--out", help="검증 산출물 디렉토리")
    ap.add_argument("--dpi", type=int, default=110)
    ap.add_argument("--no-ocr", action="store_true")
    a = ap.parse_args()
    tokens = None
    if a.tokens and os.path.exists(a.tokens):
        tokens = [ln.strip() for ln in open(a.tokens) if ln.strip()]
    ko = [k.strip() for k in a.ko.split(",")] if a.ko else None
    r = verify(a.pptx, tokens=tokens, ko=ko, out_dir=a.out, dpi=a.dpi, do_ocr=not a.no_ocr)
    c = r["counts"]
    print(f"[HARD] OOB={c['out_of_bounds']} zero={c['zero_size']} blank={c['blank_slide']} "
          f"missing={c['missing_content']}  | [ADVISORY] overlap={c['text_overlap']} "
          f"ocr_yield={c['ocr_low_yield']} ocr_ko={c['ocr_missing_korean']} (ocr_ran={r['ocr_ran']})")
    print(f"[ADVISORY/타이포] kr_neg_spacing={c['kr_neg_spacing']} "
          f"kr_wide_spacing={c['kr_wide_spacing']} kr_unsafe_font={c['kr_unsafe_font']}")
    print(f"[ADVISORY/스타일] punct_overuse={c['style_punct_overuse']} "
          f"edge_bar_overuse={c['style_edge_bar_overuse']} v_center_cram={c['style_v_center_cram']}")
    print(f"[ADVISORY/wrap] text_wrap_overflow={c['text_wrap_overflow']}")
    print(f"[ADVISORY/render] render_unavailable={c['render_unavailable']} "
          f"(렌더 도구 부재 시 OCR·이미지 빈슬라이드 검사 생략 — HARD GATE 비차단)")
    for o in r["issues"]["out_of_bounds"]:
        print(f"  [OOB] S{o['slide']}: {', '.join(o['over'])}  '{o['text']}'")
    for o in r["issues"]["ocr_low_yield"]:
        print(f"  [OCR저산출] S{o['slide']}: {o['ocr_chars']}/{o['pptx_chars']} ({o['ratio']}) — tofu/저대비 의심")
    for o in r["issues"]["kr_neg_spacing"]:
        print(f"  [타이포] S{o['slide']}: 한글 음수 자간 {o['spc_pt']}pt  '{o['text']}' — 자간 벌어짐/명조 폴백 의심")
    for o in r["issues"]["kr_wide_spacing"]:
        print(f"  [타이포] S{o['slide']}: 한글 과대 자간 {o['spc_pt']}pt  '{o['text']}'")
    for o in r["issues"]["kr_unsafe_font"]:
        print(f"  [타이포] S{o['slide']}: 한글에 비안전 폰트 '{o['font']}'  '{o['text']}' — 세리프/명조 폴백 의심")
    for o in r["issues"]["style_punct_overuse"]:
        print(f"  [스타일] S{o['slide']}: 큰 글씨 구분 부호(·/—/–) {o['count']}회 — 남발 의심, 문장으로 풀거나 레이아웃으로 분리")
    if r["issues"]["style_edge_bar_overuse"]:
        _sl = sorted({o['slide'] for o in r["issues"]["style_edge_bar_overuse"]})
        print(f"  [스타일] 좌측 액센트 바 {c['style_edge_bar_overuse']}개(S{_sl}) — 남발 의심, 강조 1~2곳만 남기기")
    for o in r["issues"]["style_v_center_cram"]:
        print(f"  [스타일] S{o['slide']}: 콘텐츠가 수직 중앙에 몰림(높이 {o['h_ratio']}H) — 3존 리듬(헤더/콘텐츠/푸터)으로 재배치")
    for o in r["issues"]["text_wrap_overflow"]:
        print(f"  [wrap] S{o['slide']}: {o['naive_lines']}문단→{o['est_lines']}줄 추정, "
              f"필요 {o['needed_in']}in > 박스 {o['box_h_in']}in  '{o['text']}' — 폭 부족 래핑 오버플로 의심(#413)")
    for o in r["issues"]["render_unavailable"]:
        print(f"  [render] {o['note']}")
    print("HARD GATE:", "PASS" if r["hard_gate_pass"] else "FAIL")
    sys.exit(0 if r["hard_gate_pass"] else 1)


if __name__ == "__main__":
    main()
