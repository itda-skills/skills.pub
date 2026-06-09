# -*- coding: utf-8 -*-
"""examples/sample/gen.py — pptx-design 스킬 동작 예제 덱 생성기.

이 스크립트는 pptx-design 스킬의 per-invocation 생성 스크립트(gen.py)가
어떤 모습이어야 하는지를 보여주는 정식 예제다.

3개 SSoT 입력:
  - content.md  : 슬라이드 콘텐츠 명세(고정 — 기본 pptx 스킬과 동일 콘텐츠 비교용)
  - data.json   : 수치 데이터(차트/표/스탯 콜아웃)
  - design.md   : 적용 DESIGN.md (여기서는 Stripe designspec — gradient mesh / 단일 indigo CTA /
                  deep navy ink / thin display + 음수 자간 / tabular numerics / pill / cream band)

설계 규칙(SPEC-PPTX-DESIGN-001 / -002 superset):
  - 도형/텍스트/표/이미지 배치는 전부 ``scripts/deckkit.py`` 공개 API 로 한다.
  - 차트 두 갈래(REQ-003): ① S3 HBM 시장 규모는 ``dk.add_native_chart`` 네이티브 편집 차트로
    시연(PowerPoint 에서 데이터 편집 가능) / ② 면적채움 추세선·끝점 라벨·슬라이더 등 고디자인
    차트는 matplotlib PNG(편집 불가). 둘 다 design.md 팔레트(기본 파랑/주황 금지).
  - gradient mesh·모티프 배경은 Pillow+numpy PNG 로 베이크(python-pptx 그라디언트 미지원 대응).
  - 한글은 deckkit.kr_font_name()(LibreOffice 또렷 고딕 — Noto Sans KR·Pretendard 우선) + 자동
    가드(REQ-001: 한글 음수 자간→0, 비안전 폰트→안전 고딕). DESIGN.md 음수 자간은 라틴 run 에만 유효.
  - 산출: examples/sample/deck.pptx + examples/sample/assets/*.png.
"""
import os
import sys
import json

# ---------------------------------------------------------------- deckkit import
HERE = os.path.dirname(os.path.abspath(__file__))
SKILL_ROOT = os.path.dirname(os.path.dirname(HERE))          # .../pptx-design
SCRIPTS = os.path.join(SKILL_ROOT, "scripts")
sys.path.insert(0, SCRIPTS)

import deckkit as dk  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

import matplotlib  # noqa: E402
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402
from PIL import Image, ImageFilter  # noqa: E402

# ---------------------------------------------------------------- paths / data
ASSETS = os.path.join(HERE, "assets")
os.makedirs(ASSETS, exist_ok=True)
DATA = json.load(open(os.path.join(HERE, "data.json"), encoding="utf-8"))
OUT = os.path.join(HERE, "deck.pptx")

# ---------------------------------------------------------------- palette (design.md = stripe.designspec.md 토큰)
PRIMARY        = "533afd"   # indigo — signature CTA
PRIMARY_DEEP   = "4434d4"
PRIMARY_PRESS  = "2e2b8c"
PRIMARY_SOFT   = "665efd"
PRIMARY_SUBDUE = "b9b9f9"
BRAND_DARK_900 = "1c1e54"   # deep navy chrome
INK            = "0d253d"   # default body text — deep navy, never pure black
INK_SECONDARY  = "273951"
INK_MUTE       = "64748d"
ON_PRIMARY     = "ffffff"
CANVAS         = "ffffff"
CANVAS_SOFT    = "f6f9fc"
CANVAS_CREAM   = "f5e9d4"
HAIRLINE       = "e3e8ee"
RUBY           = "ea2261"   # gradient/chart accent only (never a button)
MAGENTA        = "f96bee"
LEMON          = "9b6829"
SHERBET        = "f6c177"   # derived warm gradient stop
LAVENDER       = "c9c4fb"

# Fonts — Latin display(Sohne 300 근사) + 한글 안전 폰트(이 환경에서 LibreOffice 치환 안전)
F_LATIN = "Helvetica Neue"
F_KR = dk.kr_font_name()

W_IN, H_IN = dk.DEFAULT_W_IN, dk.DEFAULT_H_IN
MX = 0.85  # outer margin

prs = dk.new_deck()


def hx(h):
    return "#" + h.lstrip("#")


# ================================================================ 텍스트 헬퍼(deckkit 래퍼)
def add_paras(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.0, space_after=6, wrap=True):
    """여러 문단 텍스트 박스. paras: [[ (text, kw), ... ], ...]."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, runs in enumerate(paras):
        if i == 0:
            p = tf.paragraphs[0]
            p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
            if space_after is not None:
                p.space_after = Pt(space_after)
            for text, kw in runs:
                r = p.add_run(); r.text = text; dk.set_run_font(r, **kw)
        else:
            dk.add_paragraph(tf, runs, align=align, space_after=space_after, line_spacing=line_spacing)
    return tb


def pill_tag(slide, x, y, text, w=None, fill=PRIMARY_SUBDUE, txt=PRIMARY_DEEP, size=10.5):
    if w is None:
        w = 0.16 + len(text) * (size * 0.0105)
    h = 0.30
    sp = dk.rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.10)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    dk.set_run_font(r, name=F_LATIN, size=size, bold=False, color=txt, spacing=0.6)
    return sp


# ================================================================ Pillow: gradient mesh (design.md 시그니처)
def _blob(arr, cx, cy, rx, ry, color, strength):
    h, w, _ = arr.shape
    yy, xx = np.mgrid[0:h, 0:w]
    d = ((xx - cx) / rx) ** 2 + ((yy - cy) / ry) ** 2
    mask = np.clip(1.0 - d, 0.0, 1.0) ** 1.6
    for c in range(3):
        arr[:, :, c] = arr[:, :, c] * (1 - mask * strength) + color[c] * (mask * strength)


def _C(hex_):
    return np.array([int(hex_[i:i + 2], 16) for i in (0, 2, 4)], float)


def make_mesh(path, w=2667, h_band=900, base="ffffff"):
    """상단 1/3 을 채우는 가로 wash pastel mesh(cream→sherbet→lavender→indigo→ruby→magenta).
    하단으로 갈수록 투명 fade → 흰 캔버스로 자연스럽게 녹는다."""
    arr = np.ones((h_band, w, 3), float) * _C(base)
    _blob(arr, w * 0.06, h_band * 0.30, w * 0.30, h_band * 0.95, _C(CANVAS_CREAM), 0.95)
    _blob(arr, w * 0.20, h_band * 0.55, w * 0.26, h_band * 0.85, _C(SHERBET),      0.55)
    _blob(arr, w * 0.40, h_band * 0.20, w * 0.30, h_band * 0.90, _C(LAVENDER),     0.70)
    _blob(arr, w * 0.58, h_band * 0.50, w * 0.26, h_band * 0.95, _C(PRIMARY_SOFT), 0.62)
    _blob(arr, w * 0.74, h_band * 0.30, w * 0.24, h_band * 0.85, _C(PRIMARY),      0.50)
    _blob(arr, w * 0.88, h_band * 0.55, w * 0.24, h_band * 0.95, _C(RUBY),         0.46)
    _blob(arr, w * 0.97, h_band * 0.25, w * 0.20, h_band * 0.80, _C(MAGENTA),      0.42)
    img = Image.fromarray(arr.astype("uint8"), "RGB").filter(ImageFilter.GaussianBlur(48))
    rgba = img.convert("RGBA")
    a = np.ones((h_band, w), float)
    fade_start = int(h_band * 0.45)
    for yy in range(fade_start, h_band):
        a[yy, :] = max(0.0, 1.0 - (yy - fade_start) / (h_band - fade_start))
    rgba.putalpha(Image.fromarray((a * 255).astype("uint8"), "L"))
    rgba.save(path)
    return path


def make_mesh_strip(path, w=2667, h=210):
    """콘텐츠 슬라이드 상단 액센트용 얇은 풀블리드 mesh strip."""
    arr = np.ones((h, w, 3), float) * 255
    _blob(arr, w * 0.05, h * 0.5, w * 0.22, h * 1.6, _C(CANVAS_CREAM), 0.9)
    _blob(arr, w * 0.30, h * 0.5, w * 0.22, h * 1.6, _C(LAVENDER),     0.7)
    _blob(arr, w * 0.55, h * 0.5, w * 0.22, h * 1.6, _C(PRIMARY_SOFT), 0.6)
    _blob(arr, w * 0.78, h * 0.5, w * 0.20, h * 1.6, _C(PRIMARY),      0.5)
    _blob(arr, w * 0.95, h * 0.5, w * 0.18, h * 1.6, _C(RUBY),         0.45)
    Image.fromarray(arr.astype("uint8"), "RGB").filter(ImageFilter.GaussianBlur(36)).save(path)
    return path


def make_mesh_dark(path, w=2667, h=1500, base="0d253d"):
    """deep navy 캔버스 위 glowing indigo/ruby mesh — dark hero & 결론용."""
    arr = np.ones((h, w, 3), float) * _C(base)
    _blob(arr, w * 0.18, h * 0.18, w * 0.40, h * 0.55, _C(PRIMARY),      0.55)
    _blob(arr, w * 0.42, h * 0.10, w * 0.34, h * 0.42, _C(PRIMARY_SOFT), 0.40)
    _blob(arr, w * 0.78, h * 0.16, w * 0.36, h * 0.50, _C(RUBY),         0.30)
    _blob(arr, w * 0.92, h * 0.40, w * 0.30, h * 0.45, _C(MAGENTA),      0.22)
    _blob(arr, w * 0.10, h * 0.85, w * 0.30, h * 0.40, _C(PRIMARY_DEEP), 0.30)
    Image.fromarray(arr.astype("uint8"), "RGB").filter(ImageFilter.GaussianBlur(60)).save(path)
    return path


MESH = make_mesh(os.path.join(ASSETS, "mesh.png"))
MESH_STR = make_mesh_strip(os.path.join(ASSETS, "mesh_strip.png"))
MESH_DARK = make_mesh_dark(os.path.join(ASSETS, "mesh_dark.png"))

# ================================================================ matplotlib (design.md 팔레트 차트)
dk.mpl_korean()
plt.rcParams["axes.unicode_minus"] = False


def style_ax(ax):
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    for s in ("left", "bottom"):
        ax.spines[s].set_color(hx(HAIRLINE))
        ax.spines[s].set_linewidth(0.8)
    ax.tick_params(colors=hx(INK_MUTE), length=0, labelsize=11)
    ax.grid(axis="y", color=hx(HAIRLINE), linewidth=0.8, alpha=0.9)
    ax.set_axisbelow(True)


def chart_financials(rev, op, fname):
    yrs = DATA["financials"]["years"]
    fig, ax = plt.subplots(figsize=(6.5, 4.1), dpi=200)
    x = np.arange(len(yrs)); ww = 0.38
    b1 = ax.bar(x - ww / 2, rev, ww, label="매출", color=hx(BRAND_DARK_900), zorder=3)
    b2 = ax.bar(x + ww / 2, op, ww, label="영업이익", color=hx(PRIMARY), zorder=3)
    style_ax(ax)
    ax.set_xticks(x); ax.set_xticklabels(yrs)
    ax.axhline(0, color=hx(INK_MUTE), linewidth=0.9, zorder=2)
    lo = min(0, min(op)); hi = max(max(rev), max(op))
    ax.set_ylim(lo - hi * 0.06, hi * 1.16)
    ax.set_yticks([]); ax.grid(False); ax.spines["left"].set_visible(False)
    for b, v in zip(b1, rev):
        ax.text(b.get_x() + b.get_width() / 2, v + hi * 0.02, f"{v:g}", ha="center", va="bottom",
                color=hx(INK), fontsize=9.5)
    for b, v in zip(b2, op):
        va = "bottom" if v >= 0 else "top"
        off = hi * 0.02 if v >= 0 else -hi * 0.02
        col = hx(RUBY) if v < 0 else hx(PRIMARY_DEEP)
        ax.text(b.get_x() + b.get_width() / 2, v + off, f"{v:g}", ha="center", va=va,
                color=col, fontsize=9.5, fontweight="medium")
    ax.legend(loc="upper left", frameon=False, fontsize=11, labelcolor=hx(INK_SECONDARY))
    fig.tight_layout()
    fig.savefig(os.path.join(ASSETS, fname), transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)


def chart_price():
    pm = DATA["price_monthly"]
    labels = pm["labels"]
    s = np.array(pm["samsung"], float); h = np.array(pm["hynix"], float)
    sN = s / s[0] * 100; hN = h / h[0] * 100
    fig, ax = plt.subplots(figsize=(11.2, 3.7), dpi=200)
    xs = np.arange(len(labels))
    ax.plot(xs, hN, color=hx(PRIMARY), linewidth=2.6, zorder=4, label="SK하이닉스")
    ax.plot(xs, sN, color=hx(BRAND_DARK_900), linewidth=2.2, zorder=3, label="삼성전자")
    ax.fill_between(xs, hN, hN.min() - 5, color=hx(PRIMARY), alpha=0.06, zorder=1)
    style_ax(ax)
    ax.set_ylim(min(sN.min(), hN.min()) - 8, max(sN.max(), hN.max()) + 12)
    step = 3
    ax.set_xticks(xs[::step])
    ax.set_xticklabels([labels[i] for i in range(0, len(labels), step)], fontsize=10)
    ax.set_ylabel("상대수익률 (기준=100)", color=hx(INK_MUTE), fontsize=11)
    for arr, col, name in ((hN, PRIMARY, "하이닉스"), (sN, BRAND_DARK_900, "삼성전자")):
        ax.scatter([xs[-1]], [arr[-1]], color=hx(col), s=34, zorder=5)
        ax.text(xs[-1] + 0.3, arr[-1], f"{name} {arr[-1]:.0f}", color=hx(col),
                fontsize=10.5, va="center", fontweight="medium")
    ax.set_xlim(0, len(labels) + 3)
    ax.legend(loc="upper left", frameon=False, fontsize=11, labelcolor=hx(INK_SECONDARY), ncol=2)
    fig.tight_layout()
    fig.savefig(os.path.join(ASSETS, "s6_price.png"), transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)


def chart_share():
    d = DATA["hbm_share_pct"]
    yrs = d["years"]
    fig, ax = plt.subplots(figsize=(6.4, 4.2), dpi=200)
    x = np.arange(len(yrs)); ww = 0.5
    sk = np.array(d["sk_hynix"]); sa = np.array(d["samsung"]); mi = np.array(d["micron"])
    ax.bar(x, sk, ww, color=hx(PRIMARY), zorder=3, label="SK하이닉스")
    ax.bar(x, sa, ww, bottom=sk, color=hx(BRAND_DARK_900), zorder=3, label="삼성전자")
    ax.bar(x, mi, ww, bottom=sk + sa, color=hx(RUBY), zorder=3, label="Micron")
    style_ax(ax); ax.grid(False); ax.set_yticks([]); ax.spines["left"].set_visible(False)
    ax.set_xticks(x); ax.set_xticklabels(yrs); ax.set_ylim(0, 104)
    for i in range(len(yrs)):
        ax.text(i, sk[i] / 2, f"{sk[i]}%", ha="center", va="center", color="white", fontsize=11, fontweight="medium")
        ax.text(i, sk[i] + sa[i] / 2, f"{sa[i]}%", ha="center", va="center", color="white", fontsize=10.5)
        ax.text(i, sk[i] + sa[i] + mi[i] / 2, f"{mi[i]}%", ha="center", va="center", color="white", fontsize=10)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.07), ncol=3, frameon=False,
              fontsize=10.5, labelcolor=hx(INK_SECONDARY))
    fig.tight_layout()
    fig.savefig(os.path.join(ASSETS, "s7_share.png"), transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)


def chart_target():
    tp = DATA["target_price"]
    fig, ax = plt.subplots(figsize=(11.0, 3.7), dpi=200)
    names = [("SK하이닉스", tp["hynix"]), ("삼성전자", tp["samsung"])]
    ypos = [1.0, 0.0]
    LEFT = 0.16
    span = 1.0 - LEFT - 0.02
    for (nm, d), y in zip(names, ypos):
        bear, base, bull, cur = d["bear"], d["base"], d["bull"], d["current"]
        lo = min(bear, cur); hi = bull
        pad = (hi - lo) * 0.10
        lo -= pad; hi += pad

        def nx(v):
            return LEFT + (v - lo) / (hi - lo) * span

        xb, xba, xbu, xc = nx(bear), nx(base), nx(bull), nx(cur)
        ax.add_patch(plt.Rectangle((xb, y - 0.11), xbu - xb, 0.22, color=hx(HAIRLINE), zorder=2))
        ax.add_patch(plt.Rectangle((xb, y - 0.11), xba - xb, 0.22, color=hx(PRIMARY_SUBDUE), zorder=3))
        for xv, val, col, lab in ((xb, bear, INK_MUTE, "Bear"), (xba, base, PRIMARY, "Base"),
                                  (xbu, bull, PRIMARY_DEEP, "Bull")):
            ax.scatter([xv], [y], color=hx(col), s=52, zorder=5)
            ax.text(xv, y + 0.20, f"{lab}\n{val:,}", ha="center", va="bottom", color=hx(col),
                    fontsize=9.4, linespacing=1.15)
        ax.scatter([xc], [y], marker="D", color=hx(RUBY), s=64, zorder=6)
        ax.text(xc, y - 0.24, f"현재 {cur:,}", ha="center", va="top", color=hx(RUBY),
                fontsize=9.6, fontweight="medium")
        ax.text(LEFT - 0.02, y, nm, ha="right", va="center", color=hx(INK), fontsize=13.5)
    ax.set_ylim(-0.7, 1.7); ax.set_xlim(0, 1.0)
    for s in ("top", "right", "left", "bottom"):
        ax.spines[s].set_visible(False)
    ax.set_yticks([]); ax.set_xticks([])
    fig.tight_layout()
    fig.savefig(os.path.join(ASSETS, "s9_target.png"), transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)


chart_financials(DATA["financials"]["samsung_revenue"], DATA["financials"]["samsung_op_profit"], "s4_samsung.png")
chart_financials(DATA["financials"]["hynix_revenue"], DATA["financials"]["hynix_op_profit"], "s5_hynix.png")
chart_price()
chart_share()
chart_target()


# ================================================================ 슬라이드 공통 부품
def eyebrow(slide, x, y, text, color=PRIMARY_DEEP):
    dk.add_text(slide, x, y, 6.0, 0.3,
                [(text, dict(name=F_LATIN, size=11, bold=False, color=color, spacing=1.4))])


def footer(slide, idx, dark=False):
    col = ON_PRIMARY if dark else INK_MUTE
    dk.add_text(slide, MX, H_IN - 0.42, 9.0, 0.3,
                [("삼성전자 vs SK하이닉스 · 2026 메모리 슈퍼사이클", dict(name=F_KR, size=8.5, color=col))],
                anchor=MSO_ANCHOR.MIDDLE)
    dk.add_text(slide, W_IN - MX - 2.0, H_IN - 0.42, 2.0, 0.3,
                [(f"{idx:02d} / 11", dict(name=F_LATIN, size=8.5, color=col, spacing=0.5))],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def section_header(slide, eyebrow_txt, title_runs, y=0.70):
    slide.shapes.add_picture(MESH_STR, 0, 0, width=Inches(W_IN), height=Inches(0.55))
    dk.rect(slide, 0, 0.55, W_IN, 0.012, fill=HAIRLINE)
    eyebrow(slide, MX, y, eyebrow_txt)
    dk.add_text(slide, MX, y + 0.30, W_IN - 2 * MX, 1.0, title_runs, anchor=MSO_ANCHOR.TOP)


# ---------------------------------------------------------------- S1 표지
def s1():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS)
    s.shapes.add_picture(MESH, 0, 0, width=Inches(W_IN), height=Inches(4.3))
    dk.add_text(s, MX, 0.55, 6, 0.3,
                [("EQUITY RESEARCH", dict(name=F_LATIN, size=11, color=INK_SECONDARY, spacing=2.2))])
    pill_tag(s, W_IN - MX - 2.0, 0.52, "2026 OUTLOOK", w=2.0, fill=CANVAS, txt=PRIMARY_DEEP, size=10)
    dk.rect(s, W_IN - MX - 2.0, 0.52, 2.0, 0.30, fill=None, line=PRIMARY_SUBDUE, line_w=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_paras(s, MX, 2.30, 11.6, 2.2,
              [[("삼성전자 ", dict(name=F_KR, size=52, color=INK, spacing=-1.2)),
                ("vs", dict(name=F_LATIN, size=52, color=PRIMARY, spacing=-1.2)),
                (" SK하이닉스", dict(name=F_KR, size=52, color=INK, spacing=-1.2))]],
              line_spacing=1.02)
    dk.add_text(s, MX, 3.55, 11.6, 0.7,
                [("2026 메모리 슈퍼사이클, 누가 더 오를까", dict(name=F_KR, size=24, color=INK_SECONDARY, spacing=-0.4))])
    dk.add_text(s, MX, 4.42, 11.6, 0.5,
                [("AI · HBM 수요가 다시 그린 두 거인의 주가 지형도", dict(name=F_KR, size=15, color=INK_MUTE))])
    dk.rect(s, MX, 5.35, W_IN - 2 * MX, 0.012, fill=HAIRLINE)
    metas = [("종목코드", "005930 · 000660"), ("작성", "Equity Research (가상)"), ("기준", "2026.06")]
    cw = (W_IN - 2 * MX) / 3
    for i, (k, v) in enumerate(metas):
        x = MX + i * cw
        dk.add_text(s, x, 5.55, cw - 0.2, 0.3, [(k, dict(name=F_KR, size=10, color=INK_MUTE, spacing=0.8))])
        dk.add_text(s, x, 5.82, cw - 0.2, 0.4, [(v, dict(name=F_LATIN, size=15, color=INK, spacing=-0.3))])
    cta = dk.rect(s, MX, 6.55, 2.55, 0.46, fill=PRIMARY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = cta.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "투자 전망 살펴보기"
    dk.set_run_font(r, name=F_KR, size=13, color=ON_PRIMARY)
    dk.add_text(s, MX + 2.75, 6.55, 8.0, 0.46,
                [("※ 본 자료의 수치는 디자인 테스트용 예시입니다.", dict(name=F_KR, size=9.5, color=INK_MUTE))],
                anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------- S2 핵심 요약 (스탯 콜아웃)
def s2():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS_SOFT)
    section_header(s, "EXECUTIVE SUMMARY", [("핵심 요약", dict(name=F_KR, size=34, color=INK, spacing=-0.6))])
    stats = DATA["key_stats"]
    accents = [PRIMARY, PRIMARY_DEEP, BRAND_DARK_900, PRIMARY_DEEP]
    cw = (W_IN - 2 * MX - 3 * 0.28) / 4
    cy = 2.05; ch = 2.55
    for i, (st, ac) in enumerate(zip(stats, accents)):
        x = MX + i * (cw + 0.28)
        dk.rect(s, x, cy, cw, ch, fill=CANVAS, line=HAIRLINE, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)
        dk.rect(s, x + 0.28, cy + 0.30, 0.18, 0.18, fill=ac, shape=MSO_SHAPE.OVAL)
        dk.add_text(s, x + 0.55, cy + 0.24, cw - 0.7, 0.3,
                    [(f"0{i+1}", dict(name=F_LATIN, size=11, color=INK_MUTE, spacing=1.0))],
                    anchor=MSO_ANCHOR.MIDDLE)
        dk.add_text(s, x + 0.28, cy + 0.78, cw - 0.5, 0.9,
                    [(st["value"], dict(name=F_LATIN, size=40, color=ac, spacing=-1.4))]
                    if st["value"][0] in "$0123456789" else
                    [(st["value"], dict(name=F_KR, size=34, color=ac, spacing=-0.8))])
        dk.add_text(s, x + 0.28, cy + 1.70, cw - 0.5, 0.5,
                    [(st["label"], dict(name=F_KR, size=12.5, color=INK_SECONDARY, spacing=-0.2))], wrap=True)
        dk.add_text(s, x + 0.28, cy + 2.12, cw - 0.5, 0.3,
                    [(st["sub"], dict(name=F_KR, size=10.5, color=INK_MUTE))])
    by = 4.95
    dk.rect(s, MX, by, W_IN - 2 * MX, 1.05, fill=CANVAS_CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.035)
    dk.rect(s, MX, by, 0.10, 1.05, fill=PRIMARY)
    dk.add_text(s, MX + 0.45, by, W_IN - 2 * MX - 0.8, 1.05,
                [("AI 수요가 만든 메모리 슈퍼사이클의 두 번째 국면 ", dict(name=F_KR, size=15.5, color=INK, spacing=-0.2)),
                 ("·  ", dict(name=F_KR, size=15.5, color=PRIMARY)),
                 ("HBM 리더 ", dict(name=F_KR, size=15.5, color=INK, spacing=-0.2)),
                 ("하이닉스의 모멘텀", dict(name=F_KR, size=15.5, color=PRIMARY_DEEP, spacing=-0.2)),
                 (", 턴어라운드 폭이 큰 ", dict(name=F_KR, size=15.5, color=INK, spacing=-0.2)),
                 ("삼성전자", dict(name=F_KR, size=15.5, color=PRIMARY_DEEP, spacing=-0.2)),
                 (".", dict(name=F_KR, size=15.5, color=INK))],
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    footer(s, 2)


# ---------------------------------------------------------------- S3 시장 환경 (HBM 성장)
def s3():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS)
    section_header(s, "MARKET CONTEXT", [("왜 지금 메모리인가", dict(name=F_KR, size=34, color=INK, spacing=-0.6))])
    pts = [
        ("AI 서버·가속기 수요 폭증", "HBM(고대역폭메모리)이 메모리 업황의 새 성장 축으로 부상."),
        ("범용 DRAM/NAND 가격 반등", "재고 정상화 후 가격이 반등 사이클에 진입."),
        ("공급 타이트·협상력 상승", "캐파 제약 + 선단 공정 집중으로 가격 협상력 상승."),
    ]
    ly = 2.20
    for i, (t, d) in enumerate(pts):
        y = ly + i * 1.30
        dk.rect(s, MX, y + 0.04, 0.42, 0.42, fill=PRIMARY_SUBDUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.30)
        dk.add_text(s, MX, y + 0.04, 0.42, 0.42, [(str(i + 1), dict(name=F_LATIN, size=15, color=PRIMARY_DEEP))],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        dk.add_text(s, MX + 0.62, y, 4.7, 0.4, [(t, dict(name=F_KR, size=16, color=INK, spacing=-0.3))])
        dk.add_text(s, MX + 0.62, y + 0.40, 4.7, 0.7, [(d, dict(name=F_KR, size=12, color=INK_MUTE))], line_spacing=1.2)
    cx = 6.30; cyy = 2.05; cwd = W_IN - MX - cx; chh = 4.45
    dk.rect(s, cx, cyy, cwd, chh, fill=CANVAS_SOFT, line=HAIRLINE, line_w=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    dk.add_text(s, cx + 0.4, cyy + 0.30, cwd - 0.8, 0.4,
                [("HBM 시장 규모", dict(name=F_KR, size=15, color=INK, spacing=-0.3)),
                 ("  (USD Bn)", dict(name=F_LATIN, size=11, color=INK_MUTE))])
    dk.add_text(s, cx + 0.4, cyy + 0.70, cwd - 0.8, 0.3,
                [("2023 → 2026E · ", dict(name=F_LATIN, size=11, color=INK_MUTE, spacing=-0.2)),
                 ("7.6배 성장 · ", dict(name=F_KR, size=11, color=RUBY)),
                 ("PowerPoint 편집 가능 차트", dict(name=F_KR, size=10, color=INK_MUTE))])
    # ★REQ-003 시연: 네이티브 편집 차트(matplotlib PNG 가 아닌 진짜 pptx 차트 객체 → 수치 편집 가능)
    hbm = DATA["hbm_market_size_usd_bn"]
    dk.add_native_chart(
        s, cx + 0.34, cyy + 1.18, cwd - 0.68, chh - 1.52, "column",
        categories=[str(y) for y in hbm["years"]],
        series=[("HBM 시장 규모", hbm["values"])],
        palette=[PRIMARY], data_labels=True, number_format='"$"0.0"B"',
        legend=False, value_axis=False, gap_width=58,
        font_name=F_KR, font_color=INK_MUTE, font_size=11.5,
    )
    footer(s, 3)


# ---------------------------------------------------------------- S4/S5 기업 개요
def company_slide(idx, code_name, code, body, points, chart_png):
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS)
    section_header(s, "COMPANY PROFILE",
                   [(code_name, dict(name=F_KR, size=34, color=INK, spacing=-0.6)),
                    (f"  ({code})", dict(name=F_LATIN, size=20, color=INK_MUTE, spacing=-0.3))])
    dk.add_text(s, MX, 2.10, 5.35, 1.5, [(body, dict(name=F_KR, size=13, color=INK_SECONDARY, spacing=-0.2))],
                line_spacing=1.35)
    py = 3.55
    for i, (t, d) in enumerate(points):
        y = py + i * 0.82
        dk.rect(s, MX, y + 0.05, 0.30, 0.30, fill=PRIMARY, shape=MSO_SHAPE.OVAL)
        dk.add_text(s, MX, y + 0.05, 0.30, 0.30, [("✓", dict(name=F_LATIN, size=11, color=ON_PRIMARY))],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        dk.add_text(s, MX + 0.48, y, 4.9, 0.3, [(t, dict(name=F_KR, size=13.5, color=INK, spacing=-0.2))])
        dk.add_text(s, MX + 0.48, y + 0.34, 4.9, 0.4, [(d, dict(name=F_KR, size=10.5, color=INK_MUTE))])
    cx = 6.30; cyy = 2.05; cwd = W_IN - MX - cx; chh = 4.45
    dk.rect(s, cx, cyy, cwd, chh, fill=CANVAS_SOFT, line=HAIRLINE, line_w=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    dk.add_text(s, cx + 0.4, cyy + 0.30, cwd - 0.8, 0.4,
                [("매출 vs 영업이익 추이", dict(name=F_KR, size=15, color=INK, spacing=-0.3)),
                 ("  (조원)", dict(name=F_LATIN, size=11, color=INK_MUTE))])
    s.shapes.add_picture(chart_png, Inches(cx + 0.28), Inches(cyy + 0.95), width=Inches(cwd - 0.56))
    footer(s, idx)


def s4():
    company_slide(4, "삼성전자", "005930",
                  "메모리·파운드리·세트(스마트폰/가전)를 아우르는 종합 반도체 기업. 2024년 메모리 흑자 전환에 이어 "
                  "2025~26년 HBM 점유율 회복과 파운드리 수율 개선이 핵심 관전 포인트다.",
                  [("메모리 대규모 흑자 복귀", "다운사이클 저점 통과 후 이익 정상화"),
                   ("HBM3E 고객 인증 확대", "엔비디아 등 핵심 고객 공급망 진입 가속"),
                   ("파운드리 적자 축소", "수율 개선과 가동률 상승으로 손실 폭 축소")],
                  os.path.join(ASSETS, "s4_samsung.png"))


def s5():
    company_slide(5, "SK하이닉스", "000660",
                  "HBM 글로벌 1위. NVIDIA향 HBM3E 선점으로 메모리 업사이클의 최대 수혜주. 2023년 적자 → 2024년 흑자 전환 "
                  "→ 2026E 사상 최대 이익 경로를 밟고 있다.",
                  [("HBM 매출 비중 급상승", "고부가 HBM이 전사 이익을 견인"),
                   ("선단 DRAM(1b/1c) 리더십", "공정 기술 우위로 원가 경쟁력 확보"),
                   ("자본 효율·현금흐름 개선", "이익 급증과 함께 재무 건전성 회복")],
                  os.path.join(ASSETS, "s5_hynix.png"))


# ---------------------------------------------------------------- S6 주가 궤적 (풀폭 차트)
def s6():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS)
    section_header(s, "PRICE TRAJECTORY", [("24개월 주가 궤적", dict(name=F_KR, size=34, color=INK, spacing=-0.6))])
    dk.add_text(s, MX, 1.62, W_IN - 2 * MX, 0.4,
                [("2024 H2 메모리 조정 → 2025 회복 → 2026 동반 상승. 하이닉스가 HBM 모멘텀으로 더 가파르게 상승.",
                  dict(name=F_KR, size=12.5, color=INK_MUTE, spacing=-0.2))])
    cx = MX; cyy = 2.20; cwd = W_IN - 2 * MX; chh = 4.30
    dk.rect(s, cx, cyy, cwd, chh, fill=CANVAS_SOFT, line=HAIRLINE, line_w=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03)
    s.shapes.add_picture(os.path.join(ASSETS, "s6_price.png"), Inches(cx + 0.35), Inches(cyy + 0.45), width=Inches(cwd - 0.7))
    footer(s, 6)


# ---------------------------------------------------------------- S7 HBM 점유율
def s7():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS_SOFT)
    section_header(s, "MARKET SHARE", [("HBM, 누가 지배하는가", dict(name=F_KR, size=34, color=INK, spacing=-0.6))])
    cx = MX; cyy = 2.05; cwd = 6.0; chh = 4.45
    dk.rect(s, cx, cyy, cwd, chh, fill=CANVAS, line=HAIRLINE, line_w=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    dk.add_text(s, cx + 0.4, cyy + 0.28, cwd - 0.8, 0.4,
                [("HBM 점유율 추이", dict(name=F_KR, size=15, color=INK, spacing=-0.3)),
                 ("  (%)", dict(name=F_LATIN, size=11, color=INK_MUTE))])
    s.shapes.add_picture(os.path.join(ASSETS, "s7_share.png"), Inches(cx + 0.30), Inches(cyy + 0.85), width=Inches(cwd - 0.6))
    rx = cx + cwd + 0.45; rw = W_IN - MX - rx
    dk.add_text(s, rx, 2.30, rw, 0.4, [("선두는 SK하이닉스", dict(name=F_KR, size=18, color=INK, spacing=-0.4))])
    dk.add_text(s, rx, 2.85, rw, 1.0, [("50%", dict(name=F_LATIN, size=64, color=PRIMARY, spacing=-2.0))])
    dk.add_text(s, rx, 3.95, rw, 0.4, [("2026E SK하이닉스 HBM 점유율 · 글로벌 1위", dict(name=F_KR, size=11.5, color=INK_MUTE))])
    dk.rect(s, rx, 4.45, rw, 0.012, fill=HAIRLINE)
    dk.add_text(s, rx, 4.65, rw, 1.6,
                [("SK하이닉스가 50%대 선두를 유지하는 가운데 삼성이 추격하고 Micron의 점유율이 상승 중이다. ",
                  dict(name=F_KR, size=12.5, color=INK_SECONDARY, spacing=-0.2)),
                 ("점유율은 곧 AI 수혜 강도", dict(name=F_KR, size=12.5, color=PRIMARY_DEEP)),
                 ("를 의미한다.", dict(name=F_KR, size=12.5, color=INK_SECONDARY))],
                line_spacing=1.35)
    footer(s, 7)


# ---------------------------------------------------------------- S8 밸류에이션 표 (도형으로 직접 그림)
def s8():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS)
    section_header(s, "VALUATION", [("비싼가, 싼가", dict(name=F_KR, size=34, color=INK, spacing=-0.6))])
    val = DATA["valuation_2026E"]
    metrics = val["metrics"]
    cols = [("지표", None, None), ("삼성전자", val["samsung"], BRAND_DARK_900),
            ("SK하이닉스", val["hynix"], PRIMARY), ("Micron (peer)", val["micron_peer"], INK_MUTE)]
    tx = MX; ty = 2.20
    tw = W_IN - 2 * MX
    col0 = 3.0
    colw = (tw - col0) / 3
    rowh = 0.86
    headh = 0.86
    nrow = len(metrics)
    feat_j = 2
    feat_x = tx + col0 + (feat_j - 1) * colw
    feat_top = ty
    feat_h = headh + nrow * rowh
    feat_bot = feat_top + feat_h

    for i in range(nrow):
        if i % 2 == 1:
            ry = ty + headh + i * rowh
            dk.rect(s, tx, ry, tw, rowh, fill=CANVAS_SOFT)
    dk.rect(s, feat_x, feat_top, colw, feat_h, fill=BRAND_DARK_900,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)

    dk.add_text(s, tx, ty, col0, headh, [("2026E 기준", dict(name=F_KR, size=12, color=INK_MUTE, spacing=-0.2))],
                anchor=MSO_ANCHOR.MIDDLE)
    for j in range(1, 4):
        x = tx + col0 + (j - 1) * colw
        nm, _, ac = cols[j]
        if nm == "SK하이닉스":
            pill_tag(s, x + colw / 2 - 0.65, ty + 0.16, "BEST VALUE", w=1.30, fill=PRIMARY, txt=ON_PRIMARY, size=8.5)
            dk.add_text(s, x, ty + 0.48, colw, 0.34, [(nm, dict(name=F_KR, size=14, color=ON_PRIMARY, spacing=-0.3))],
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            dk.add_text(s, x, ty, colw, headh, [(nm, dict(name=F_KR, size=14, color=ac, spacing=-0.3))],
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    def hairline_row(yv, thick):
        dk.rect(s, tx, yv, feat_x - tx, thick, fill=HAIRLINE)
        rx = feat_x + colw
        dk.rect(s, rx, yv, (tx + tw) - rx, thick, fill=HAIRLINE)

    hairline_row(ty + headh, 0.014)
    for i in range(nrow):
        hairline_row(ty + headh + (i + 1) * rowh, 0.008)

    for i, m in enumerate(metrics):
        ry = ty + headh + i * rowh
        dk.add_text(s, tx + 0.15, ry, col0 - 0.2, rowh, [(m, dict(name=F_KR, size=13, color=INK, spacing=-0.2))],
                    anchor=MSO_ANCHOR.MIDDLE)
        for j in range(1, 4):
            x = tx + col0 + (j - 1) * colw
            nm, vals, ac = cols[j]
            v = vals[i]
            txtcol = ON_PRIMARY if nm == "SK하이닉스" else INK
            dk.add_text(s, x, ry, colw, rowh, [(f"{v:g}", dict(name=F_LATIN, size=18, color=txtcol, spacing=-0.6))],
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    by = feat_bot + 0.30
    dk.rect(s, MX, by, tw, 0.85, fill=CANVAS_CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
    dk.rect(s, MX, by, 0.10, 0.85, fill=RUBY)
    dk.add_text(s, MX + 0.45, by, tw - 0.8, 0.85,
                [("이익 정점 논쟁이 멀티플 확장을 제약 ", dict(name=F_KR, size=13, color=INK, spacing=-0.2)),
                 ("·  ", dict(name=F_KR, size=13, color=RUBY)),
                 ("멀티플보다 이익 지속성이 관건", dict(name=F_KR, size=13, color=PRIMARY_DEEP)),
                 ("이다.", dict(name=F_KR, size=13, color=INK))],
                anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 8)


# ---------------------------------------------------------------- S9 목표주가 시나리오
def s9():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS)
    section_header(s, "PRICE TARGETS", [("목표주가 시나리오", dict(name=F_KR, size=34, color=INK, spacing=-0.6))])
    dk.add_text(s, MX, 1.62, W_IN - 2 * MX, 0.4,
                [("Bear / Base / Bull 3구간 ", dict(name=F_KR, size=12.5, color=INK_MUTE, spacing=-0.2)),
                 ("·  ", dict(name=F_KR, size=12.5, color=PRIMARY)),
                 ("현재가(◆) 대비 상승여력을 표시.", dict(name=F_KR, size=12.5, color=INK_MUTE, spacing=-0.2))])
    cx = MX; cyy = 2.10; cwd = W_IN - 2 * MX; chh = 3.55
    dk.rect(s, cx, cyy, cwd, chh, fill=CANVAS_SOFT, line=HAIRLINE, line_w=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.03)
    s.shapes.add_picture(os.path.join(ASSETS, "s9_target.png"), Inches(cx + 0.35), Inches(cyy + 0.35), width=Inches(cwd - 0.7))
    by = 5.95
    items = [("삼성전자 컨센서스 평균", "108,000", "원", BRAND_DARK_900),
             ("SK하이닉스 컨센서스 평균", "390,000", "원", PRIMARY)]
    iw = (W_IN - 2 * MX - 0.3) / 2
    for i, (lab, num, unit, ac) in enumerate(items):
        x = MX + i * (iw + 0.3)
        dk.rect(s, x, by, iw, 0.95, fill=CANVAS, line=HAIRLINE, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        dk.rect(s, x, by, 0.09, 0.95, fill=ac)
        dk.add_text(s, x + 0.35, by + 0.16, iw - 0.5, 0.3, [(lab, dict(name=F_KR, size=11, color=INK_MUTE))])
        dk.add_text(s, x + 0.35, by + 0.44, iw - 0.5, 0.4,
                    [(num, dict(name=F_LATIN, size=24, color=ac, spacing=-0.8)),
                     (" " + unit, dict(name=F_KR, size=13, color=INK_MUTE))])
    footer(s, 9)


# ---------------------------------------------------------------- S10 리스크 그리드 (3x2)
def s10():
    s = dk.blank_slide(prs); dk.set_bg(s, CANVAS_SOFT)
    section_header(s, "RISK FACTORS", [("무엇이 시나리오를 깨뜨리나", dict(name=F_KR, size=34, color=INK, spacing=-0.6))])
    risks = DATA["risks"]
    cols, rows = 3, 2
    gx = 0.26; gy = 0.26
    gw = (W_IN - 2 * MX - (cols - 1) * gx) / cols
    gh = 1.92
    sy = 2.05
    accents = [PRIMARY, PRIMARY_DEEP, RUBY, BRAND_DARK_900, MAGENTA, PRIMARY_SOFT]
    for i, r in enumerate(risks):
        cc = i % cols; rr = i // cols
        x = MX + cc * (gw + gx); y = sy + rr * (gh + gy)
        dk.rect(s, x, y, gw, gh, fill=CANVAS, line=HAIRLINE, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        ac = accents[i]
        dk.rect(s, x + 0.30, y + 0.30, 0.50, 0.50, fill=ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.28)
        dk.add_text(s, x + 0.30, y + 0.30, 0.50, 0.50, [(f"{i+1:02d}", dict(name=F_LATIN, size=14, color=ON_PRIMARY))],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        dk.add_text(s, x + 0.95, y + 0.32, gw - 1.15, 0.5, [(r["title"], dict(name=F_KR, size=14.5, color=INK, spacing=-0.3))],
                    anchor=MSO_ANCHOR.MIDDLE)
        dk.add_text(s, x + 0.30, y + 0.95, gw - 0.58, 0.85, [(r["desc"], dict(name=F_KR, size=11, color=INK_MUTE))],
                    line_spacing=1.25)
    footer(s, 10)


# ---------------------------------------------------------------- S11 결론 (dark navy + mesh)
def s11():
    s = dk.blank_slide(prs); dk.set_bg(s, INK)
    s.shapes.add_picture(MESH_DARK, 0, 0, width=Inches(W_IN), height=Inches(H_IN))
    eyebrow(s, MX, 0.62, "CONCLUSION", color=PRIMARY_SOFT)
    dk.add_text(s, MX, 0.92, 9.0, 0.9, [("결론 & 투자의견", dict(name=F_KR, size=36, color=ON_PRIMARY, spacing=-0.8))])
    cons = DATA["analyst_consensus"]
    cards = [
        ("SK하이닉스", "강력 매수", cons["hynix"], PRIMARY, "HBM 모멘텀의 최대 수혜, 이익 가시성 최상."),
        ("삼성전자", "매수", cons["samsung"], PRIMARY_SOFT, "턴어라운드 폭과 자산가치 매력, HBM 추격 성공이 리레이팅 트리거."),
    ]
    cy = 2.15; ch = 2.55
    cw = (W_IN - 2 * MX - 0.40) / 2
    for i, (nm, rating, c, ac, note) in enumerate(cards):
        x = MX + i * (cw + 0.40)
        dk.rect(s, x, cy, cw, ch, fill=BRAND_DARK_900, line=PRIMARY_PRESS, line_w=1.0,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        dk.add_text(s, x + 0.45, cy + 0.32, cw - 0.9, 0.4, [(nm, dict(name=F_KR, size=20, color=ON_PRIMARY, spacing=-0.4))])
        pill_tag(s, x + 0.45, cy + 0.85, rating, w=0.55 + len(rating) * 0.22, fill=ac, txt=ON_PRIMARY, size=12)
        bhs = f"Buy {c['buy']}  ·  Hold {c['hold']}  ·  Sell {c['sell']}"
        dk.add_text(s, x + 0.45, cy + 1.40, cw - 0.9, 0.4, [(bhs, dict(name=F_LATIN, size=14, color="b9c4d4", spacing=-0.2))])
        dk.add_text(s, x + 0.45, cy + 1.85, cw - 0.9, 0.6, [(note, dict(name=F_KR, size=11.5, color="aeb9cc"))], line_spacing=1.25)
    by = 5.05
    dk.rect(s, MX, by, W_IN - 2 * MX, 0.95, fill="14315a", shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    dk.rect(s, MX, by, 0.10, 0.95, fill=RUBY)
    dk.add_text(s, MX + 0.45, by, W_IN - 2 * MX - 0.8, 0.95,
                [("성장 모멘텀 = 하이닉스, 안정·복원력 = 삼성. ", dict(name=F_KR, size=13.5, color=ON_PRIMARY, spacing=-0.2)),
                 ("사이클 정점 신호(메모리 가격 피크아웃) 모니터링", dict(name=F_KR, size=13.5, color=MAGENTA)),
                 ("이 핵심.", dict(name=F_KR, size=13.5, color=ON_PRIMARY))],
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    dk.add_text(s, MX, 6.30, W_IN - 2 * MX, 0.4,
                [("본 자료의 모든 수치는 디자인 테스트용 예시이며 투자 권유가 아닙니다.",
                  dict(name=F_KR, size=9.5, color="aab6cc"))])
    dk.add_text(s, W_IN - MX - 2.0, 6.30, 2.0, 0.4, [("11 / 11", dict(name=F_LATIN, size=9.5, color="aab6cc", spacing=0.5))],
                align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- build
s1(); s2(); s3(); s4(); s5(); s6(); s7(); s8(); s9(); s10(); s11()
dk.save_deck(prs, OUT)
print("saved", OUT, "| 11 slides | F_KR =", F_KR)
